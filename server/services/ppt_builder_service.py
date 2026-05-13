"""PPT 스타일 빌더 서비스 (M3)

generated_pptx_styled 컬렉션의 structured JSON 을 읽어 python-pptx 로 실제 .pptx
파일을 빌드한다. 슬라이드마다 SSE 이벤트를 송출하기 위해 async generator 로
구현되어 있다.

참고 문서: 데스크탑 PPT생성스.md 의 다음 절을 기준으로 좌표/색/폰트를 결정.
  - 3절 (디자인 토큰)
  - 4.3절 (배경 그라데이션)
  - 5절 (12개 패턴 라이브러리)
  - 6절 (헬퍼 함수: 페이지 인디케이터 / contentHeader / chapterDivider / 빅스탯)
  - 7절 (트러블슛 5건)

핵심 5개 패턴 우선 구현: cover / toc / chapter / content_3col / closing.
나머지 7개 패턴은 build_fallback() 으로 단순 텍스트 슬라이드 형태로 강등 처리.

좌표 단위: pptxgenjs (Node) 의 inch 값을 python-pptx 의 Inches() 로 동일하게 사용.
캔버스: 10" x 5.625" (16:9 LAYOUT_16x9).
"""

from __future__ import annotations

import asyncio
import os
import shutil
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Optional

from bson import ObjectId
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from config import settings
from services.mongo_service import get_db
from services import ppt_asset_service
from services import ppt_style_service
from services.pptx_import_service import parse_pptx_to_slides


# ============ 캔버스 상수 ============

CANVAS_W_IN = 10.0
CANVAS_H_IN = 5.625

# 기본 폰트 (windows 기준 Malgun Gothic)
DEFAULT_TITLE_FONT = "Malgun Gothic"
DEFAULT_BODY_FONT = "Malgun Gothic"


# ============ 색상 / 폰트 헬퍼 ============

def _hex_to_rgbcolor(hex_str: str, fallback: str = "#000000") -> RGBColor:
    """'#RRGGBB' → RGBColor. 잘못된 입력 시 fallback 사용."""
    s = hex_str if isinstance(hex_str, str) and hex_str else fallback
    s = s.strip().lstrip("#")
    if len(s) != 6:
        s = fallback.lstrip("#")
    try:
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
    except Exception:
        r, g, b = 0, 0, 0
    return RGBColor(r, g, b)


def _color(tokens: dict, key: str, fallback: str = "#000000") -> RGBColor:
    """design_tokens.colors[key] → RGBColor (없으면 fallback)"""
    colors = (tokens.get("colors") or {}) if isinstance(tokens, dict) else {}
    return _hex_to_rgbcolor(colors.get(key, fallback), fallback)


def _color_hex(tokens: dict, key: str, fallback: str = "#000000") -> str:
    """design_tokens.colors[key] → '#RRGGBB' (없으면 fallback)"""
    colors = (tokens.get("colors") or {}) if isinstance(tokens, dict) else {}
    val = colors.get(key)
    if isinstance(val, str) and val.strip():
        return val if val.startswith("#") else "#" + val
    return fallback


def _font_family(skill_file: dict, role: str = "title") -> str:
    """skill_file.fonts 에서 폰트 family 추출.

    role: "title" | "body" → fonts 배열에서 매칭되는 가장 첫 번째 폰트.
    skill_file.design_tokens.fonts.title_font_id / body_font_id 와 fonts 배열의
    _id 를 매칭한다. 없으면 시스템 기본 폰트(Malgun Gothic) 사용.
    """
    if not isinstance(skill_file, dict):
        return DEFAULT_TITLE_FONT
    design_tokens = skill_file.get("design_tokens") or {}
    fonts_meta = (design_tokens.get("fonts") or {})
    fonts_list = skill_file.get("fonts") or []

    target_id = fonts_meta.get(f"{role}_font_id")
    if target_id and fonts_list:
        for f in fonts_list:
            if str(f.get("_id")) == str(target_id):
                family = (f.get("family") or f.get("name") or "").strip()
                if family:
                    return family
    # role 우선 → 첫 폰트 fallback
    if fonts_list:
        first = fonts_list[0]
        family = (first.get("family") or first.get("name") or "").strip()
        if family:
            return family
    return DEFAULT_TITLE_FONT if role == "title" else DEFAULT_BODY_FONT


def _font_size(tokens: dict, key: str, fallback: int) -> int:
    """design_tokens.fonts.sizes[key] → int Pt 값"""
    if not isinstance(tokens, dict):
        return fallback
    sizes = (tokens.get("fonts") or {}).get("sizes") or {}
    val = sizes.get(key)
    if isinstance(val, (int, float)) and val > 0:
        return int(val)
    return fallback


# ============ 텍스트 박스 헬퍼 ============

def _add_textbox(
    slide,
    text: str,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    font_family: str = DEFAULT_TITLE_FONT,
    font_size: int = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = RGBColor(0, 0, 0),
    align: str = "left",
    anchor: str = "top",
    line_spacing: Optional[float] = None,
    char_spacing: Optional[int] = None,
):
    """간단한 텍스트 박스 추가 헬퍼. 단일 paragraph + 단일 run."""
    box = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # 첫 paragraph 활용
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    if line_spacing:
        p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text or ""
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    try:
        run.font.color.rgb = color
    except Exception:
        pass
    if char_spacing is not None:
        # python-pptx 는 charSpacing 을 직접 지원하지 않으므로 OXML 로 주입
        try:
            from pptx.oxml.ns import qn
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))  # 100 단위
        except Exception:
            pass
    return box


def _add_rect(
    slide,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    fill_color: Optional[RGBColor] = None,
    line_color: Optional[RGBColor] = None,
    line_width_pt: float = 0.0,
    shape_type=MSO_SHAPE.RECTANGLE,
):
    """도형(기본 사각형) 추가. fill_color None → 무채움, line_color None → 무외곽선."""
    # 너비/높이 0 인 경우 최소값 보장
    w = max(w_in, 0.001)
    h = max(h_in, 0.001)
    shape = slide.shapes.add_shape(shape_type, Inches(x_in), Inches(y_in), Inches(w), Inches(h))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        try:
            shape.fill.fore_color.rgb = fill_color
        except Exception:
            pass
    if line_color is None or line_width_pt <= 0:
        shape.line.fill.background()
    else:
        try:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(line_width_pt)
        except Exception:
            pass
    # 도형 텍스트 비활성 (placeholder 텍스트 노출 방지)
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


# ============ 트러블슛 7.2 — 본문 텍스트 길이 자동 축약 ============

def _truncate_desc(text: str, max_chars: int = 90) -> str:
    """카드 본문이 max_chars 초과 시 말줄임 처리. .md 7.2 룰."""
    if not isinstance(text, str):
        return ""
    s = text.strip()
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 1].rstrip() + "…"


# ============ 6절 — 헬퍼 함수 ============

def add_page_indicator(slide, page_num: int, dark: bool, tokens: dict):
    """우하단 페이지 번호 + 짧은 가로 라인 (.md 6.1)"""
    color = RGBColor(0xFF, 0xFF, 0xFF) if dark else _color(tokens, "primary", "#1C60EF")
    _add_textbox(
        slide, str(page_num),
        x_in=8.9, y_in=5.2, w_in=0.7, h_in=0.3,
        font_size=10, bold=True, color=color, align="right",
    )
    _add_rect(slide, 9.6, 5.34, 0.3, 0.015, fill_color=color)


def add_content_header(slide, label: str, title: str, tokens: dict,
                       title_font: str = DEFAULT_TITLE_FONT,
                       body_font: str = DEFAULT_BODY_FONT):
    """상단 라벨 칩(좌측 컬러 바 + 라벨 텍스트) + 본문 제목 (.md 6.2)"""
    primary = _color(tokens, "primary", "#1C60EF")
    ink = _color(tokens, "ink", "#0B1E3F")

    # 좌측 컬러 바 (직사각형)
    _add_rect(slide, 0.55, 0.62, 0.32, 0.04, fill_color=primary)

    # 라벨 텍스트
    _add_textbox(
        slide, label or "",
        x_in=0.95, y_in=0.47, w_in=5.0, h_in=0.3,
        font_family=body_font, font_size=10, color=primary,
        char_spacing=6,
    )

    # 타이틀
    _add_textbox(
        slide, title or "",
        x_in=0.5, y_in=0.82, w_in=9.0, h_in=0.7,
        font_family=title_font, font_size=26, bold=True, color=ink,
    )


def add_chapter_divider_content(slide, chapter: str, ch_title: str, ch_sub: str,
                                 page_num: int, part_total: int, tokens: dict,
                                 title_font: str = DEFAULT_TITLE_FONT,
                                 body_font: str = DEFAULT_BODY_FONT):
    """거대 챕터 번호 + 가로 디바이더 + 제목 + 서브타이틀 + PART x/total 칩 (.md 6.3)

    트러블슛 7.3 대응: 챕터 번호 fontSize 240 → 180 으로 다운.
    """
    primary = _color(tokens, "primary", "#1C60EF")
    ink = _color(tokens, "ink", "#0B1E3F")
    grey = _color(tokens, "grey", "#6B7B9A")
    light = _color(tokens, "light", "#DBE8FE")

    # CHAPTER 라벨
    _add_textbox(
        slide, "CHAPTER",
        x_in=0.55, y_in=0.62, w_in=3.0, h_in=0.3,
        font_family=body_font, font_size=10, color=primary, char_spacing=6,
    )

    # 거대 챕터 번호 (180pt 고정 - 트러블슛 7.3)
    _add_textbox(
        slide, chapter or "",
        x_in=0.4, y_in=0.9, w_in=5.0, h_in=2.4,
        font_family=title_font, font_size=180, bold=True, color=primary,
        anchor="top",
    )

    # 가로 디바이더
    _add_rect(slide, 0.6, 3.65, 0.8, 0.05, fill_color=primary)

    # 챕터 타이틀
    _add_textbox(
        slide, ch_title or "",
        x_in=0.55, y_in=3.82, w_in=9.0, h_in=0.6,
        font_family=title_font, font_size=30, bold=True, color=ink,
    )

    # 서브타이틀
    if ch_sub:
        _add_textbox(
            slide, ch_sub,
            x_in=0.55, y_in=4.45, w_in=9.0, h_in=0.5,
            font_family=body_font, font_size=13, color=grey,
        )

    # PART x/total 칩 (우측 상단 인근)
    chip_x, chip_y, chip_w, chip_h = 7.4, 2.65, 2.0, 0.75
    _add_rect(slide, chip_x, chip_y, chip_w, chip_h,
              fill_color=light, line_color=primary, line_width_pt=0.5)
    _add_textbox(
        slide, f"PART {chapter}/{part_total:02d}",
        x_in=chip_x, y_in=chip_y, w_in=chip_w, h_in=chip_h,
        font_family=body_font, font_size=14, bold=True, color=primary,
        align="center", anchor="middle", char_spacing=4,
    )

    # 페이지 인디케이터
    add_page_indicator(slide, page_num, dark=False, tokens=tokens)


def split_stat_text(slide, value: str, unit: str,
                    x_in: float, y_in: float, w_in: float, h_in: float,
                    tokens: dict, title_font: str = DEFAULT_TITLE_FONT):
    """빅 스탯: 큰 숫자(value) + 작은 단위(unit) 을 별개 텍스트 박스로 분리.

    .md 7.1 트러블슛 — 단일 박스에 혼합 사이즈를 넣으면 wrap 발생. 분리 처리.
    """
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light = _color(tokens, "light", "#DBE8FE")
    # 숫자
    _add_textbox(
        slide, value or "",
        x_in=x_in, y_in=y_in, w_in=w_in * 0.78, h_in=h_in,
        font_family=title_font, font_size=88, bold=True, color=white,
        anchor="middle",
    )
    # 단위 (오른쪽으로 인접)
    _add_textbox(
        slide, unit or "",
        x_in=x_in + w_in * 0.78, y_in=y_in, w_in=w_in * 0.22, h_in=h_in,
        font_family=title_font, font_size=48, bold=True, color=light,
        anchor="middle",
    )


# ============ 배경 이미지 부착 ============

def _add_background_image(slide, prs, bg_path: str):
    """슬라이드 풀블리드 배경 이미지 부착."""
    if not bg_path or not os.path.exists(bg_path):
        return
    try:
        slide.shapes.add_picture(
            bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    except Exception as e:
        print(f"[Builder] 배경 이미지 부착 실패 {bg_path}: {e}")


# ============ 5개 핵심 패턴 빌더 ============

def build_cover(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """표지 (.md 5절 cover)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    white = RGBColor(0xFF, 0xFF, 0xFF)
    light = _color(design_tokens, "light", "#DBE8FE")

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    description = slide_data.get("description", "")
    presenter = slide_data.get("presenter", "")

    # 좌측 큰 텍스트 영역 (전체 화이트)
    _add_textbox(
        slide, title,
        x_in=0.6, y_in=1.6, w_in=8.8, h_in=1.6,
        font_family=title_font, font_size=56, bold=True, color=white,
    )

    if subtitle:
        _add_textbox(
            slide, subtitle,
            x_in=0.6, y_in=3.3, w_in=8.8, h_in=0.6,
            font_family=title_font, font_size=22, bold=True, color=light,
        )

    if description:
        _add_textbox(
            slide, _truncate_desc(description, 120),
            x_in=0.6, y_in=4.0, w_in=8.8, h_in=0.8,
            font_family=body_font, font_size=13, color=white,
            line_spacing=1.3,
        )

    if presenter:
        _add_textbox(
            slide, presenter,
            x_in=0.6, y_in=4.95, w_in=8.8, h_in=0.4,
            font_family=body_font, font_size=11, color=light, char_spacing=4,
        )

    # 장식 원 (.md 5절 cover - 우상단 외곽선 원)
    primary = _color(design_tokens, "primary", "#1C60EF")
    _add_rect(slide, 8.3, 0.4, 1.3, 1.3, fill_color=None,
              line_color=light, line_width_pt=0.75, shape_type=MSO_SHAPE.OVAL)


def build_toc(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """목차 - 좌측 헤더 + 우측 챕터 리스트 카드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    primary = _color(design_tokens, "primary", "#1C60EF")
    ink = _color(design_tokens, "ink", "#0B1E3F")
    grey = _color(design_tokens, "grey", "#6B7B9A")
    line = _color(design_tokens, "line", "#D5E2FB")
    white = RGBColor(0xFF, 0xFF, 0xFF)

    # 좌측 헤더
    _add_textbox(
        slide, "CONTENTS",
        x_in=0.55, y_in=0.62, w_in=3.0, h_in=0.3,
        font_family=body_font, font_size=10, color=primary, char_spacing=6,
    )
    _add_rect(slide, 0.55, 0.97, 0.32, 0.04, fill_color=primary)
    _add_textbox(
        slide, slide_data.get("title", "목차"),
        x_in=0.5, y_in=1.1, w_in=4.3, h_in=1.4,
        font_family=title_font, font_size=36, bold=True, color=ink,
    )
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, subtitle,
            x_in=0.55, y_in=2.55, w_in=4.3, h_in=2.0,
            font_family=body_font, font_size=12, color=grey, line_spacing=1.4,
        )

    # 우측 챕터 리스트 카드
    items = slide_data.get("items") or slide_data.get("chapters") or []
    if not items:
        # cards 필드 fallback
        cards = slide_data.get("cards") or []
        items = [{"no": c.get("no", str(i+1).zfill(2)),
                  "title": c.get("title", ""),
                  "desc": c.get("desc", "")} for i, c in enumerate(cards)]

    max_cards = 5
    items = items[:max_cards]
    n = max(1, len(items))
    card_x = 5.1
    card_w = 4.4
    top_y = 0.55
    bottom_y = 5.05
    available = bottom_y - top_y
    gap = 0.13
    card_h = max(0.5, (available - gap * (n - 1)) / n)

    for i, it in enumerate(items):
        y = top_y + i * (card_h + gap)
        _add_rect(slide, card_x, y, card_w, card_h,
                  fill_color=white, line_color=line, line_width_pt=0.5)
        no = it.get("no") if isinstance(it, dict) else str(i + 1).zfill(2)
        title = it.get("title") if isinstance(it, dict) else str(it)
        desc = it.get("desc") if isinstance(it, dict) else ""
        # 좌측 번호
        _add_textbox(
            slide, str(no),
            x_in=card_x + 0.18, y_in=y + 0.10, w_in=0.7, h_in=card_h - 0.2,
            font_family=title_font, font_size=22, bold=True, color=primary,
            anchor="middle",
        )
        # 타이틀
        _add_textbox(
            slide, title or "",
            x_in=card_x + 0.95, y_in=y + 0.12, w_in=card_w - 1.05, h_in=0.4,
            font_family=title_font, font_size=14, bold=True, color=ink,
            anchor="middle",
        )
        if desc:
            _add_textbox(
                slide, _truncate_desc(desc, 60),
                x_in=card_x + 0.95, y_in=y + card_h - 0.4, w_in=card_w - 1.05, h_in=0.3,
                font_family=body_font, font_size=10, color=grey,
            )

    add_page_indicator(slide, slide_data.get("index", 2), dark=False, tokens=design_tokens)


def build_chapter(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """챕터 디바이더 (.md 5절 chapter)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    chapter = slide_data.get("chapter") or slide_data.get("no") or "01"
    ch_title = slide_data.get("title", "")
    ch_sub = slide_data.get("subtitle", "") or slide_data.get("description", "")
    page_num = slide_data.get("index", 0)
    part_total = slide_data.get("part_total") or slide_data.get("total_parts") or 5

    add_chapter_divider_content(
        slide, str(chapter), ch_title, ch_sub, page_num, int(part_total),
        design_tokens, title_font=title_font, body_font=body_font,
    )


def build_content_3col(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """3컬럼 카드 (.md 5절 content_3col)

    번호(01·02·03) + 우상단 아이콘 자리 + 중앙 정렬 텍스트 카드 3장.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    primary = _color(design_tokens, "primary", "#1C60EF")
    ink = _color(design_tokens, "ink", "#0B1E3F")
    grey = _color(design_tokens, "grey", "#6B7B9A")
    line = _color(design_tokens, "line", "#D5E2FB")
    white = RGBColor(0xFF, 0xFF, 0xFF)
    primary_hex = _color_hex(design_tokens, "primary", "#1C60EF")

    label = slide_data.get("label", "")
    title = slide_data.get("title", "")
    add_content_header(slide, label, title, design_tokens,
                       title_font=title_font, body_font=body_font)

    cards = slide_data.get("cards") or []
    cards = cards[:3]
    if len(cards) < 3:
        # 부족 시 더미로 채워 레이아웃 균형 유지
        while len(cards) < 3:
            cards.append({"no": str(len(cards) + 1).zfill(2), "title": "", "desc": "", "icon": ""})

    # 본문 콘텐츠 영역 y: 1.70 ~ 5.05 (.md 3.3)
    top_y = 1.70
    card_h = 3.35
    gap = 0.15
    total_w = 9.0
    card_w = (total_w - gap * 2) / 3.0
    left_x = 0.5

    # 아이콘 사이즈 (이미지 임시 폴더에 PNG 저장 후 add_picture)
    icon_dir = Path(settings.UPLOAD_DIR) / "ppt_styled" / "_icons"
    icon_dir.mkdir(parents=True, exist_ok=True)

    for i, c in enumerate(cards):
        x = left_x + i * (card_w + gap)
        # 카드 박스
        _add_rect(slide, x, top_y, card_w, card_h,
                  fill_color=white, line_color=line, line_width_pt=0.5)

        # 좌상단 큰 번호
        no = c.get("no") or str(i + 1).zfill(2)
        _add_textbox(
            slide, str(no),
            x_in=x + 0.25, y_in=top_y + 0.20, w_in=1.4, h_in=0.7,
            font_family=title_font, font_size=36, bold=True, color=primary,
        )

        # 우상단 아이콘 자리 (PNG 임베드)
        icon_key = (c.get("icon") or "").strip()
        if icon_key:
            try:
                icon_bytes = ppt_asset_service.icon_placeholder(icon_key, primary_hex, size=256)
                icon_path = icon_dir / f"icon_{icon_key}_{primary_hex.lstrip('#')}.png"
                if not icon_path.exists():
                    with open(icon_path, "wb") as f:
                        f.write(icon_bytes)
                slide.shapes.add_picture(
                    str(icon_path),
                    Inches(x + card_w - 0.85), Inches(top_y + 0.2),
                    Inches(0.6), Inches(0.6),
                )
            except Exception as e:
                print(f"[Builder] 아이콘 임베드 실패 {icon_key}: {e}")

        # 카드 타이틀 (중앙 정렬)
        card_title = c.get("title") or ""
        _add_textbox(
            slide, card_title,
            x_in=x + 0.25, y_in=top_y + 1.15, w_in=card_w - 0.5, h_in=0.9,
            font_family=title_font, font_size=16, bold=True, color=ink,
            align="center", anchor="middle", line_spacing=1.2,
        )

        # 카드 본문 (60~90자 축약)
        desc = _truncate_desc(c.get("desc") or "", 90)
        _add_textbox(
            slide, desc,
            x_in=x + 0.25, y_in=top_y + 2.1, w_in=card_w - 0.5, h_in=1.1,
            font_family=body_font, font_size=11, color=grey,
            align="center", line_spacing=1.4,
        )

    add_page_indicator(slide, slide_data.get("index", 0),
                       dark=False, tokens=design_tokens)


def build_closing(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """마감 슬라이드 (.md 5절 closing)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    white = RGBColor(0xFF, 0xFF, 0xFF)
    light = _color(design_tokens, "light", "#DBE8FE")

    thanks = slide_data.get("title", "Thank You")
    summary = slide_data.get("summary") or slide_data.get("description", "")
    presenter = slide_data.get("presenter", "")

    _add_textbox(
        slide, thanks,
        x_in=0.6, y_in=1.6, w_in=8.8, h_in=1.6,
        font_family=title_font, font_size=64, bold=True, color=white,
    )

    if summary:
        _add_textbox(
            slide, _truncate_desc(summary, 200),
            x_in=0.6, y_in=3.4, w_in=8.8, h_in=1.4,
            font_family=body_font, font_size=14, color=light,
            line_spacing=1.4,
        )

    if presenter:
        _add_textbox(
            slide, presenter,
            x_in=0.6, y_in=5.0, w_in=8.8, h_in=0.4,
            font_family=body_font, font_size=11, color=light, char_spacing=4,
        )


def build_fallback(prs, slide_data: dict, skill_file: dict, bg_path: str):
    """5개 핵심 패턴 외 → content_3col 강등 또는 단순 텍스트 슬라이드.

    cards 가 있고 길이가 2~4 사이면 content_3col 로 강등.
    그 외에는 헤더 + 단일 본문 단순 슬라이드.
    """
    cards = slide_data.get("cards") or []
    if 2 <= len(cards) <= 4:
        # cards 를 정확히 3개로 정규화 (잘라내거나 빈 카드 추가)
        normalized = list(cards[:3])
        while len(normalized) < 3:
            normalized.append({"no": str(len(normalized) + 1).zfill(2),
                               "title": "", "desc": "", "icon": ""})
        new_data = dict(slide_data)
        new_data["cards"] = normalized
        build_content_3col(prs, new_data, skill_file, bg_path)
        return

    # 단순 텍스트 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background_image(slide, prs, bg_path)

    design_tokens = skill_file.get("design_tokens") or {}
    title_font = _font_family(skill_file, "title")
    body_font = _font_family(skill_file, "body")

    label = slide_data.get("label", "")
    title = slide_data.get("title", "")
    add_content_header(slide, label, title, design_tokens,
                       title_font=title_font, body_font=body_font)

    desc = slide_data.get("description") or slide_data.get("body") or ""
    points = slide_data.get("points") or []
    if points and isinstance(points, list):
        desc = desc + ("\n\n" if desc else "") + "\n".join(f"• {p}" for p in points)

    grey = _color(design_tokens, "grey", "#6B7B9A")
    _add_textbox(
        slide, desc or "(내용이 비어 있습니다)",
        x_in=0.5, y_in=1.70, w_in=9.0, h_in=3.35,
        font_family=body_font, font_size=13, color=grey,
        line_spacing=1.5,
    )

    add_page_indicator(slide, slide_data.get("index", 0),
                       dark=False, tokens=design_tokens)


# 패턴 ID → 빌더 함수 매핑
_BUILDER_REGISTRY = {
    "cover": build_cover,
    "toc": build_toc,
    "chapter": build_chapter,
    "content_3col": build_content_3col,
    "closing": build_closing,
}


def _pick_builder(template_id: str):
    """template_id 에 맞는 빌더 함수 반환. 미정의 시 fallback."""
    return _BUILDER_REGISTRY.get((template_id or "").strip(), build_fallback)


# ============ 메인 빌더 (스트리밍) ============

async def build_pptx_styled_stream(
    project_id: str,
    style_id: Optional[str] = None,
    output_dir: Optional[str] = None,
):
    """구조화 JSON 을 PPTX 로 빌드하면서 SSE 이벤트를 yield.

    Yields tuples of (event_type, payload):
      - ("slide_start", {"index", "total", "template_id", "title"})
      - ("bg_generated", {"index"})
      - ("slide_done",  {"index", "total"})
      - ("result", {"pptx_path", "pptx_url", "total_slides"})
      - ("error", {"message"})
    """
    db = get_db()

    # 1) structured 로드
    rec = await db.generated_pptx_styled.find_one({"project_id": project_id})
    if not rec or not rec.get("structured"):
        yield ("error", {"message": "구조화된 슬라이드 데이터가 없습니다."})
        return

    structured = rec["structured"] or {}
    slides = structured.get("slides") or []
    if not slides:
        yield ("error", {"message": "구조화된 슬라이드 목록이 비어 있습니다."})
        return

    # 2) style_id 확정 (param > record)
    sid = style_id or rec.get("style_id")
    if not sid:
        yield ("error", {"message": "style_id 가 없습니다."})
        return

    skill_file = await ppt_style_service.generate_skill_file(sid)
    if not skill_file:
        yield ("error", {"message": "스타일 스킬 파일을 로드할 수 없습니다."})
        return

    design_tokens = skill_file.get("design_tokens") or {}

    # 3) 출력 경로
    base_dir = Path(output_dir) if output_dir else (
        Path(settings.UPLOAD_DIR) / "ppt_styled" / str(project_id)
    )
    base_dir.mkdir(parents=True, exist_ok=True)

    # 4) Presentation 객체 (16:9)
    prs = Presentation()
    prs.slide_width = Inches(CANVAS_W_IN)
    prs.slide_height = Inches(CANVAS_H_IN)

    total = len(slides)

    # 5) 슬라이드 루프
    for i, slide_data in enumerate(slides, start=1):
        if not isinstance(slide_data, dict):
            slide_data = {}
        slide_data.setdefault("index", i)
        template_id = slide_data.get("template_id") or "content_3col"

        yield ("slide_start", {
            "index": i,
            "total": total,
            "template_id": template_id,
            "title": slide_data.get("title", ""),
        })

        # 5a) 배경 JPG 생성
        bg_params = ppt_asset_service.slide_background_origin(i, template_id)
        # 토큰 → hex 치환
        c_in_token = bg_params.pop("color_in_token", "light")
        c_out_token = bg_params.pop("color_out_token", "white")
        bg_kwargs = dict(bg_params)
        bg_kwargs["color_in"] = _color_hex(design_tokens, c_in_token,
                                            "#DBE8FE" if c_in_token == "light" else "#1C60EF")
        bg_kwargs["color_out"] = _color_hex(design_tokens, c_out_token,
                                             "#FFFFFF" if c_out_token == "white" else "#071949")
        bg_path = base_dir / f"bg_{i:02d}.jpg"
        try:
            await ppt_asset_service.generate_background(
                str(bg_path),
                canvas_size=(1920, 1080),
                **bg_kwargs,
            )
            yield ("bg_generated", {"index": i})
        except Exception as e:
            print(f"[Builder] 배경 생성 실패 idx={i}: {e}")
            # 실패해도 진행 (배경 없이 빌드)
            bg_path = ""

        # 5b) 패턴별 빌더 호출
        try:
            builder = _pick_builder(template_id)
            builder(prs, slide_data, skill_file, str(bg_path) if bg_path else "")
        except Exception as e:
            print(f"[Builder] 슬라이드 빌드 실패 idx={i} pattern={template_id}: {e}")
            # 빌드 실패 시 fallback 으로 재시도
            try:
                build_fallback(prs, slide_data, skill_file, str(bg_path) if bg_path else "")
            except Exception as e2:
                print(f"[Builder] fallback 도 실패 idx={i}: {e2}")

        yield ("slide_done", {"index": i, "total": total})

        # 너무 빠른 루프에서 이벤트 송출 보장
        await asyncio.sleep(0)

    # 6) PPTX 저장
    pptx_path = base_dir / "presentation.pptx"
    try:
        prs.save(str(pptx_path))
    except Exception as e:
        yield ("error", {"message": f"PPTX 저장 실패: {e}"})
        return

    # 6a) M16.A — 빌드 후 후처리 (폰트 임베드 등)
    post_result = {"embedded_fonts": 0, "errors": []}
    try:
        post_result = _post_process_pptx(str(pptx_path), skill_file)
        if post_result.get("embedded_fonts"):
            yield ("fonts_embedded", {"count": post_result["embedded_fonts"]})
    except Exception as e:
        print(f"[Builder] post_process 전체 실패: {e}")

    # 7) DB 업데이트 (빌드 완료)
    pptx_url = f"/uploads/ppt_styled/{project_id}/presentation.pptx"
    await db.generated_pptx_styled.update_one(
        {"project_id": project_id},
        {"$set": {
            "status": "built",
            "pptx_path": str(pptx_path),
            "pptx_url": pptx_url,
            "embedded_fonts_count": post_result.get("embedded_fonts", 0),
            "built_at": datetime.utcnow(),
            "updated_at": datetime.utcnow(),
        }},
    )

    # 8) M5 — 빌드된 PPTX 를 파싱하여 generated_slides 컬렉션에 저장
    #    (사용자측 기존 미리보기/편집 렌더러가 이 데이터를 그대로 소비)
    parsed_slide_count = 0
    parsed_object_count = 0
    parse_error_msg: Optional[str] = None
    try:
        parse_result = await asyncio.to_thread(parse_pptx_to_slides, str(pptx_path))
        parsed_slides = parse_result.get("slides") or []

        # 기존 도큐먼트 제거 (정합성 보장 — 빌드는 idempotent)
        await db.generated_slides.delete_many({"project_id": project_id})

        now = datetime.utcnow()
        slide_docs = []
        for idx, parsed in enumerate(parsed_slides):
            if not isinstance(parsed, dict):
                continue
            objects = parsed.get("objects") or []
            slide_doc = {
                "project_id": project_id,
                "order": idx + 1,
                "objects": objects,
                "slide_meta": parsed.get("slide_meta") or {},
                "background_image": parsed.get("background_image"),
                "items": [],
                "created_at": now,
                "updated_at": now,
            }
            slide_docs.append(slide_doc)
            parsed_object_count += len(objects)

        if slide_docs:
            await db.generated_slides.insert_many(slide_docs)
        parsed_slide_count = len(slide_docs)

        # generated_pptx_styled 도큐먼트에 파싱 메타 업데이트
        await db.generated_pptx_styled.update_one(
            {"project_id": project_id},
            {"$set": {
                "parsed_slide_count": parsed_slide_count,
                "parsed_object_count": parsed_object_count,
                "parsed_at": datetime.utcnow(),
                "updated_at": datetime.utcnow(),
            }, "$unset": {"parse_error": ""}},
        )

        yield ("parsed", {
            "slide_count": parsed_slide_count,
            "object_count": parsed_object_count,
        })
    except Exception as e:
        parse_error_msg = str(e)
        print(f"[Builder] PPTX 파싱/저장 실패: {parse_error_msg}")
        try:
            await db.generated_pptx_styled.update_one(
                {"project_id": project_id},
                {"$set": {
                    "parse_error": parse_error_msg,
                    "updated_at": datetime.utcnow(),
                }},
            )
        except Exception as e2:
            print(f"[Builder] parse_error DB 기록 실패: {e2}")
        yield ("parse_failed", {"message": parse_error_msg})

    yield ("result", {
        "pptx_path": str(pptx_path),
        "pptx_url": pptx_url,
        "total_slides": total,
        "parsed_slide_count": parsed_slide_count,
        "parsed_object_count": parsed_object_count,
    })


__all__ = [
    "build_pptx_styled_stream",
    "build_pptx_styled_from_designs_stream",
    "render_design_spec",
    "build_cover", "build_toc", "build_chapter",
    "build_content_3col", "build_closing", "build_fallback",
    "add_page_indicator", "add_content_header",
    "add_chapter_divider_content", "split_stat_text",
    # M16.A — OOXML 헬퍼
    "_apply_run_text_alpha", "_apply_shape_fill_alpha",
    "_apply_gradient_fill", "_set_run_char_spacing",
    "_post_process_pptx",
]


# ============ 자유 디자인 스펙 렌더러 (재설계 흐름) ============

# 30개 아이콘 키 (디자이너 프롬프트와 동기 유지)
_ICON_KEYS_VALID = {
    "users", "atom", "ship", "bomb", "shield", "fire", "chart", "star",
    "lightning", "target", "brain", "gear", "lock", "leaf", "globe", "rocket",
    "eye", "check", "warning", "idea", "clock", "money", "building", "network",
    "book", "cloud", "code", "database", "mobile", "palette",
}


def _resolve_url_to_local(url: str) -> Optional[str]:
    """`/uploads/...` URL → settings.UPLOAD_DIR 기준 절대경로. 외부/실패 시 None."""
    if not url or not isinstance(url, str):
        return None
    if not url.startswith("/uploads/"):
        return None
    rel = url[len("/uploads/"):]
    full = os.path.join(settings.UPLOAD_DIR, rel.replace("/", os.sep))
    return full if os.path.isfile(full) else None


def _safe_float(value, default: float) -> float:
    try:
        v = float(value)
        if v != v:  # NaN
            return default
        return v
    except (TypeError, ValueError):
        return default


def _safe_int(value, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _clamp_rect(x: float, y: float, w: float, h: float) -> tuple:
    """캔버스(0,0)~(CANVAS_W,CANVAS_H) 내로 clamp. 음수/0 너비는 최소값 보장."""
    x = max(0.0, min(x, CANVAS_W_IN - 0.01))
    y = max(0.0, min(y, CANVAS_H_IN - 0.01))
    w = max(0.01, min(w, CANVAS_W_IN - x))
    h = max(0.01, min(h, CANVAS_H_IN - y))
    return x, y, w, h


# ============ M16.A — OOXML 깊은 조작 헬퍼 ============

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _alpha_val(opacity: float) -> str:
    """opacity 0.0~1.0 → OOXML alpha val (0~100000 1/1000 단위 문자열)."""
    op = max(0.0, min(1.0, float(opacity)))
    return str(int(op * 100000))


def _apply_solid_fill_with_alpha(fill_xml, hex_color: str, opacity: float) -> None:
    """fill 의 a:solidFill > a:srgbClr 자식에 a:alpha 를 보장.

    fill_xml: pptx 의 fill XML element (보통 shape.fill._xPr)
    """
    if fill_xml is None:
        return
    sf = fill_xml.find(qn("a:solidFill"))
    if sf is None:
        # solidFill 노드가 없으면 추가
        sf = etree.SubElement(fill_xml, qn("a:solidFill"))
    # 기존 자식 제거 후 srgbClr 추가 (일관성 유지)
    for child in list(sf):
        sf.remove(child)
    srgb = etree.SubElement(sf, qn("a:srgbClr"))
    srgb.set("val", hex_color.lstrip("#").upper())
    if 0.0 <= opacity < 1.0:
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", _alpha_val(opacity))


def _apply_shape_fill_alpha(shape, opacity: float) -> None:
    """이미 solid fill 이 설정된 shape 에 alpha 만 주입 (색 변경 없음)."""
    if not (0.0 <= opacity < 1.0):
        return
    try:
        fill_xml = shape.fill._xPr
        sf = fill_xml.find(qn("a:solidFill"))
        if sf is None:
            return
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is None:
            return
        # 기존 alpha 제거 후 재추가
        for ch in list(srgb):
            if ch.tag == qn("a:alpha"):
                srgb.remove(ch)
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", _alpha_val(opacity))
    except Exception as e:
        print(f"[OOXML] shape alpha 실패: {e}")


def _apply_run_text_alpha(run, opacity: float) -> None:
    """text run 의 솔리드 fill 에 alpha 주입 → 텍스트 자체 투명도."""
    if not (0.0 <= opacity < 1.0):
        return
    try:
        rPr = run._r.get_or_add_rPr()
        # 기존 solidFill 찾기, 없으면 추가
        sf = rPr.find(qn("a:solidFill"))
        if sf is None:
            # run.font.color 가 설정돼 있어도 OOXML 상으로는 a:solidFill 이 아닌
            # a:fill 계열일 수 있다. 그 경우 새 solidFill 을 추가
            sf = etree.SubElement(rPr, qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is None:
            # run.font.color.rgb 값에서 hex 추출 시도
            try:
                rgb = run.font.color.rgb
                hex_val = str(rgb)
            except Exception:
                hex_val = "000000"
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.set("val", hex_val.upper())
        # 기존 alpha 제거 후 재추가
        for ch in list(srgb):
            if ch.tag == qn("a:alpha"):
                srgb.remove(ch)
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", _alpha_val(opacity))
    except Exception as e:
        print(f"[OOXML] run alpha 실패: {e}")


def _set_run_char_spacing(run, spacing_pt: float) -> None:
    """text run 에 a:rPr/@spc 설정 (1/100 pt 단위).

    spacing_pt 0.0 이면 노옵.
    """
    if not spacing_pt:
        return
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(float(spacing_pt) * 100)))
    except Exception as e:
        print(f"[OOXML] char spacing 실패: {e}")


def _apply_gradient_fill(shape, color_from_hex: str, color_to_hex: str,
                          angle_deg: float = 0.0,
                          opacity_from: float = 1.0,
                          opacity_to: float = 1.0) -> None:
    """shape 에 선형 그라데이션 fill 적용 (OOXML a:gradFill).

    python-pptx 0.6.x 는 gradFill 직접 지원이 부족하므로 spPr 의 fill 자식을 교체.
    angle_deg: 0=좌→우, 90=상→하 (OOXML 의 a:lin angle 단위는 1/60000 deg).
    """
    try:
        spPr = shape.fill._xPr  # spPr 또는 그 후손
        # 기존 fill 자식 모두 제거 (solidFill, noFill, gradFill, blipFill, pattFill)
        FILL_TAGS = {qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                     qn("a:blipFill"), qn("a:pattFill"), qn("a:grpFill")}
        for ch in list(spPr):
            if ch.tag in FILL_TAGS:
                spPr.remove(ch)

        # gradFill 노드 추가
        gradFill = etree.SubElement(spPr, qn("a:gradFill"))
        gradFill.set("flip", "none")
        gradFill.set("rotWithShape", "1")

        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

        # gs0 (start)
        gs0 = etree.SubElement(gsLst, qn("a:gs"))
        gs0.set("pos", "0")
        srgb0 = etree.SubElement(gs0, qn("a:srgbClr"))
        srgb0.set("val", color_from_hex.lstrip("#").upper())
        if 0.0 <= opacity_from < 1.0:
            a0 = etree.SubElement(srgb0, qn("a:alpha"))
            a0.set("val", _alpha_val(opacity_from))

        # gs1 (end)
        gs1 = etree.SubElement(gsLst, qn("a:gs"))
        gs1.set("pos", "100000")
        srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
        srgb1.set("val", color_to_hex.lstrip("#").upper())
        if 0.0 <= opacity_to < 1.0:
            a1 = etree.SubElement(srgb1, qn("a:alpha"))
            a1.set("val", _alpha_val(opacity_to))

        # lin (direction)
        lin = etree.SubElement(gradFill, qn("a:lin"))
        # OOXML lin angle 은 1/60000 deg, 0 = 좌→우
        lin.set("ang", str(int(float(angle_deg) * 60000) % (360 * 60000)))
        lin.set("scaled", "0")

        # tileRect 빈 자식 (mandatory in some renderers)
        etree.SubElement(gradFill, qn("a:tileRect"))
    except Exception as e:
        print(f"[OOXML] gradient fill 실패: {e}")


# ============ M16.A — 새 region 타입 ============

def _render_ghost_text_region(slide, region: dict, skill_file: dict):
    """거대 반투명 배경 텍스트 (장식용)."""
    x = _safe_float(region.get("x"), 0.5)
    y = _safe_float(region.get("y"), 0.5)
    w = _safe_float(region.get("w"), 5.0)
    h = _safe_float(region.get("h"), 4.0)
    x, y, w, h = _clamp_rect(x, y, w, h)

    text = str(region.get("text", "") or "")
    font_family = (region.get("font_family") or "").strip() \
        or _font_family(skill_file, "title")
    font_size = _safe_int(region.get("font_size"), 200)
    bold = bool(region.get("bold", True))
    italic = bool(region.get("italic", False))
    color_hex = region.get("color") or "#1C60EF"
    if not isinstance(color_hex, str) or not color_hex.strip():
        color_hex = "#1C60EF"
    opacity = _safe_float(region.get("opacity"), 0.08)
    align = (region.get("align") or "left").lower()
    if align not in ("left", "center", "right"):
        align = "left"
    valign = (region.get("valign") or "top").lower()
    anchor = {"top": "top", "middle": "middle", "bottom": "bottom"}.get(valign, "top")

    box = _add_textbox(
        slide, text,
        x_in=x, y_in=y, w_in=w, h_in=h,
        font_family=font_family,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=_hex_to_rgbcolor(color_hex, "#1C60EF"),
        align=align,
        anchor=anchor,
    )
    # 첫 run 에 alpha 주입
    try:
        run = box.text_frame.paragraphs[0].runs[0]
        _apply_run_text_alpha(run, opacity)
    except Exception:
        pass


def _render_page_indicator_region(slide, region: dict, skill_file: dict):
    """페이지 인디케이터: 작은 숫자 + 가는 가로 라인."""
    x = _safe_float(region.get("x"), 8.5)
    y = _safe_float(region.get("y"), 5.2)
    w = _safe_float(region.get("w"), 1.2)
    h = _safe_float(region.get("h"), 0.3)
    x, y, w, h = _clamp_rect(x, y, w, h)

    number = str(region.get("number", "") or "")
    total = region.get("total")
    color_hex = region.get("color") or "#1C60EF"
    if not isinstance(color_hex, str) or not color_hex.strip():
        color_hex = "#1C60EF"
    color = _hex_to_rgbcolor(color_hex, "#1C60EF")
    font_family = (region.get("font_family") or "").strip() \
        or _font_family(skill_file, "body")
    font_size = _safe_int(region.get("font_size"), 9)

    # 숫자 (좌측)
    text = number
    if total is not None:
        try:
            text = f"{number} / {int(total):02d}"
        except (TypeError, ValueError):
            pass
    num_w = w * 0.55
    _add_textbox(
        slide, text,
        x_in=x, y_in=y, w_in=num_w, h_in=h,
        font_family=font_family, font_size=font_size, bold=True,
        color=color, align="right", anchor="middle", char_spacing=3,
    )
    # 우측 짧은 라인
    line_y = y + h / 2 - 0.008
    line_x = x + num_w + 0.06
    line_w = max(0.1, w - num_w - 0.06)
    _add_rect(slide, line_x, line_y, line_w, 0.015, fill_color=color)


def _render_chip_region(slide, region: dict, skill_file: dict):
    """라벨 chip: 보더 박스 + 중앙 텍스트."""
    x = _safe_float(region.get("x"), 8.0)
    y = _safe_float(region.get("y"), 2.5)
    w = _safe_float(region.get("w"), 1.4)
    h = _safe_float(region.get("h"), 0.4)
    x, y, w, h = _clamp_rect(x, y, w, h)

    text = str(region.get("text", "") or "")
    color_hex = region.get("color") or "#1C60EF"
    if not isinstance(color_hex, str) or not color_hex.strip():
        color_hex = "#1C60EF"
    color = _hex_to_rgbcolor(color_hex, "#1C60EF")

    fill_raw = region.get("fill") or "none"
    fill_color: Optional[RGBColor] = None
    fill_opacity = _safe_float(region.get("fill_opacity"), 1.0)
    if isinstance(fill_raw, str) and fill_raw and fill_raw.lower() != "none":
        fill_color = _hex_to_rgbcolor(fill_raw, "#FFFFFF")

    stroke_width = _safe_float(region.get("stroke_width"), 0.75)

    font_family = (region.get("font_family") or "").strip() \
        or _font_family(skill_file, "body")
    font_size = _safe_int(region.get("font_size"), 10)
    bold = bool(region.get("bold", True))

    align = (region.get("align") or "center").lower()
    if align not in ("left", "center", "right"):
        align = "center"
    valign = (region.get("valign") or "middle").lower()
    anchor = {"top": "top", "middle": "middle", "bottom": "bottom"}.get(valign, "middle")

    char_spacing_raw = region.get("char_spacing")
    try:
        char_spacing = int(char_spacing_raw) if char_spacing_raw is not None else 2
    except (TypeError, ValueError):
        char_spacing = 2

    # 보더 박스
    shape = _add_rect(
        slide, x, y, w, h,
        fill_color=fill_color,
        line_color=color,
        line_width_pt=stroke_width,
    )
    # fill opacity (fill 색이 있을 때만)
    if fill_color is not None and 0.0 <= fill_opacity < 1.0:
        _apply_shape_fill_alpha(shape, fill_opacity)

    # 중앙 텍스트
    _add_textbox(
        slide, text,
        x_in=x, y_in=y, w_in=w, h_in=h,
        font_family=font_family, font_size=font_size, bold=bold,
        color=color, align=align, anchor=anchor, char_spacing=char_spacing,
    )


def _render_accent_line_region(slide, region: dict):
    """카드/섹션 상단 강조용 짧은 라인 (얇은 직사각형 또는 connector)."""
    x = _safe_float(region.get("x"), 0.5)
    y = _safe_float(region.get("y"), 1.7)
    w = _safe_float(region.get("w"), 3.0)
    h = _safe_float(region.get("h"), 0.05)
    x, y, w, h = _clamp_rect(x, y, w, h)

    color_hex = region.get("color") or "#1C60EF"
    if not isinstance(color_hex, str) or not color_hex.strip():
        color_hex = "#1C60EF"
    color = _hex_to_rgbcolor(color_hex, "#1C60EF")
    opacity = _safe_float(region.get("opacity"), 1.0)

    orientation = (region.get("orientation") or "horizontal").lower()
    if orientation == "vertical" and h < w:
        # 세로 라인이지만 사용자가 w/h 를 잘못 줬을 때 → 보정
        w, h = h, w

    shape = _add_rect(slide, x, y, w, h, fill_color=color)
    if 0.0 <= opacity < 1.0:
        _apply_shape_fill_alpha(shape, opacity)


def _render_decoration_set_region(slide, region: dict):
    """여러 도형(원/사각형)을 한꺼번에 렌더. 내부 shapes 배열을 순회."""
    shapes_list = region.get("shapes") or []
    if not isinstance(shapes_list, list):
        return
    for sub in shapes_list:
        if not isinstance(sub, dict):
            continue
        # type 키가 없는 sub-shape 도 shape region 으로 취급
        try:
            _render_shape_region(slide, sub)
        except Exception as e:
            print(f"[Builder-Spec] decoration_set sub-shape 실패: {e}")


# ============ M16.A — gradient 도형 (shape region 확장) ============

def _render_gradient_shape_region(slide, region: dict):
    """gradient fill 을 가진 shape region.

    region 필수 필드: gradient_from, gradient_to, (gradient_angle)
    """
    x = _safe_float(region.get("x"), 0.0)
    y = _safe_float(region.get("y"), 0.0)
    w = _safe_float(region.get("w"), 1.0)
    h = _safe_float(region.get("h"), 1.0)
    x, y, w, h = _clamp_rect(x, y, w, h)

    shape_kind = (region.get("shape") or "rectangle").lower()
    shape_type = MSO_SHAPE.RECTANGLE if shape_kind != "ellipse" else MSO_SHAPE.OVAL

    g_from = region.get("gradient_from") or "#1C60EF"
    g_to = region.get("gradient_to") or "#DBE8FE"
    g_angle = _safe_float(region.get("gradient_angle"), 0.0)
    op_from = _safe_float(region.get("gradient_opacity_from"), 1.0)
    op_to = _safe_float(region.get("gradient_opacity_to"), 1.0)

    stroke_raw = region.get("stroke")
    stroke_width = _safe_float(region.get("stroke_width"), 0.0)
    stroke_color: Optional[RGBColor] = None
    if isinstance(stroke_raw, str) and stroke_raw and stroke_raw.lower() != "none":
        stroke_color = _hex_to_rgbcolor(stroke_raw, "#000000")

    # 일단 solid fill 로 만든 뒤 fill XML 을 gradient 로 교체
    shape = _add_rect(
        slide, x, y, w, h,
        fill_color=_hex_to_rgbcolor(g_from, "#FFFFFF"),
        line_color=stroke_color,
        line_width_pt=stroke_width if stroke_color is not None else 0.0,
        shape_type=shape_type,
    )
    _apply_gradient_fill(shape, g_from, g_to, g_angle, op_from, op_to)


def _render_text_region(slide, region: dict, skill_file: dict):
    """text region 렌더."""
    x = _safe_float(region.get("x"), 0.5)
    y = _safe_float(region.get("y"), 0.5)
    w = _safe_float(region.get("w"), 9.0)
    h = _safe_float(region.get("h"), 0.5)
    x, y, w, h = _clamp_rect(x, y, w, h)

    text = region.get("text", "")
    if not isinstance(text, str):
        text = str(text or "")

    # font_family: 빈 문자열이면 skill_file 첫 폰트 또는 기본
    font_family = (region.get("font_family") or "").strip()
    if not font_family:
        font_family = _font_family(skill_file, "body")

    font_size = _safe_int(region.get("font_size"), 14)
    bold = bool(region.get("bold", False))
    italic = bool(region.get("italic", False))
    color_hex = region.get("color") or "#000000"
    if not isinstance(color_hex, str):
        color_hex = "#000000"
    color = _hex_to_rgbcolor(color_hex, "#000000")

    align = (region.get("align") or "left").lower()
    if align not in ("left", "center", "right"):
        align = "left"
    valign = (region.get("valign") or "top").lower()
    anchor = {"top": "top", "middle": "middle", "bottom": "bottom"}.get(valign, "top")

    cs_raw = region.get("char_spacing")
    char_spacing = None
    if cs_raw is not None:
        try:
            char_spacing = int(cs_raw)
        except (TypeError, ValueError):
            char_spacing = None

    opacity = _safe_float(region.get("opacity"), 1.0)

    box = _add_textbox(
        slide, text,
        x_in=x, y_in=y, w_in=w, h_in=h,
        font_family=font_family,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=color,
        align=align,
        anchor=anchor,
        char_spacing=char_spacing,
    )
    # 텍스트 자체 opacity (run alpha)
    if 0.0 <= opacity < 1.0:
        try:
            run = box.text_frame.paragraphs[0].runs[0]
            _apply_run_text_alpha(run, opacity)
        except Exception:
            pass


def _render_shape_region(slide, region: dict):
    """shape region 렌더 (rectangle | ellipse | line).

    M16.A: gradient fill / opacity / stroke opacity 지원.
    region 에 gradient_from/gradient_to 가 있으면 gradient 우선 적용.
    """
    x = _safe_float(region.get("x"), 0.0)
    y = _safe_float(region.get("y"), 0.0)
    w = _safe_float(region.get("w"), 1.0)
    h = _safe_float(region.get("h"), 1.0)
    x, y, w, h = _clamp_rect(x, y, w, h)

    shape_kind = (region.get("shape") or "rectangle").lower()

    fill_raw = region.get("fill")
    stroke_raw = region.get("stroke")
    stroke_width = _safe_float(region.get("stroke_width"), 1.0)
    opacity = _safe_float(region.get("opacity"), 1.0)

    # gradient 우선 처리
    g_from = region.get("gradient_from")
    g_to = region.get("gradient_to")
    if isinstance(g_from, str) and g_from and isinstance(g_to, str) and g_to \
            and shape_kind != "line":
        _render_gradient_shape_region(slide, region)
        return

    fill_color: Optional[RGBColor] = None
    if isinstance(fill_raw, str) and fill_raw and fill_raw.lower() != "none":
        fill_color = _hex_to_rgbcolor(fill_raw, "#FFFFFF")

    stroke_color: Optional[RGBColor] = None
    if isinstance(stroke_raw, str) and stroke_raw and stroke_raw.lower() != "none":
        stroke_color = _hex_to_rgbcolor(stroke_raw, "#000000")

    if shape_kind == "line":
        # 연결선 (Line) 으로 처리. add_connector(MSO_CONNECTOR.STRAIGHT, ...)
        try:
            from pptx.enum.shapes import MSO_CONNECTOR
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x), Inches(y),
                Inches(x + w), Inches(y + h),
            )
            if stroke_color is not None and stroke_width > 0:
                conn.line.color.rgb = stroke_color
                conn.line.width = Pt(stroke_width)
        except Exception as e:
            print(f"[Builder-Spec] line shape 실패: {e}")
        return

    shape_type = MSO_SHAPE.RECTANGLE if shape_kind != "ellipse" else MSO_SHAPE.OVAL

    shape = _add_rect(
        slide,
        x, y, w, h,
        fill_color=fill_color,
        line_color=stroke_color,
        line_width_pt=stroke_width if stroke_color is not None else 0.0,
        shape_type=shape_type,
    )

    # opacity (fill 만 적용. 1.0 이외일 때만 OOXML 으로 alpha 주입)
    if fill_color is not None and 0.0 <= opacity < 1.0:
        _apply_shape_fill_alpha(shape, opacity)


def _render_icon_region(slide, region: dict, icon_dir: Path):
    """icon region 렌더 (PIL 로 PNG 생성 후 add_picture)."""
    x = _safe_float(region.get("x"), 0.0)
    y = _safe_float(region.get("y"), 0.0)
    w = _safe_float(region.get("w"), 0.6)
    h = _safe_float(region.get("h"), 0.6)
    x, y, w, h = _clamp_rect(x, y, w, h)

    icon_key = (region.get("icon_key") or "").strip()
    if icon_key and icon_key not in _ICON_KEYS_VALID:
        # 잘못된 키는 빈 원으로 그림
        icon_key = ""
    color_hex = region.get("color") or "#1C60EF"
    if not isinstance(color_hex, str) or not color_hex.strip():
        color_hex = "#1C60EF"

    try:
        icon_bytes = ppt_asset_service.icon_placeholder(icon_key, color_hex, size=256)
        icon_dir.mkdir(parents=True, exist_ok=True)
        # 캐시 파일명에 색상 hex 포함
        safe_color = color_hex.lstrip("#")
        icon_path = icon_dir / f"icon_{icon_key or 'blank'}_{safe_color}.png"
        if not icon_path.exists():
            with open(icon_path, "wb") as f:
                f.write(icon_bytes)
        slide.shapes.add_picture(
            str(icon_path),
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
    except Exception as e:
        print(f"[Builder-Spec] icon 렌더 실패 {icon_key}: {e}")


def _render_image_region(slide, region: dict):
    """image region 렌더. image_url 이 /uploads/ 로 시작하면 로컬 변환 후 add_picture."""
    x = _safe_float(region.get("x"), 0.0)
    y = _safe_float(region.get("y"), 0.0)
    w = _safe_float(region.get("w"), 4.0)
    h = _safe_float(region.get("h"), 3.0)
    x, y, w, h = _clamp_rect(x, y, w, h)

    url = region.get("image_url") or ""
    local = _resolve_url_to_local(url)
    if not local:
        # 외부 URL 또는 미존재 → 건너뜀
        return
    try:
        slide.shapes.add_picture(
            local,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
    except Exception as e:
        print(f"[Builder-Spec] image 렌더 실패 {url}: {e}")


def _render_background(slide, prs, design_spec: dict, asset_dir: Path, slide_index: int):
    """design_spec.background 에 따라 슬라이드 배경 렌더."""
    bg = design_spec.get("background") or {}
    if not isinstance(bg, dict):
        return

    btype = (bg.get("type") or "none").lower()

    if btype == "image":
        url = bg.get("image_url") or ""
        local = _resolve_url_to_local(url)
        if local:
            try:
                slide.shapes.add_picture(
                    local, 0, 0,
                    width=prs.slide_width, height=prs.slide_height,
                )
            except Exception as e:
                print(f"[Builder-Spec] 배경 이미지 부착 실패: {e}")
        return

    if btype == "solid":
        color_hex = bg.get("from_color") or bg.get("color") or "#FFFFFF"
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                prs.slide_width, prs.slide_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgbcolor(color_hex, "#FFFFFF")
            shape.line.fill.background()
            # 배경은 맨 뒤로 — python-pptx 는 z-order 직접 제어가 제한적이므로
            # 가장 먼저 추가했기 때문에 자동으로 뒤에 위치
        except Exception as e:
            print(f"[Builder-Spec] 단색 배경 실패: {e}")
        return

    if btype == "gradient":
        style = (bg.get("style") or "radial").lower()
        if style not in ("radial", "linear"):
            style = "radial"
        from_color = bg.get("from_color") or "#DBE8FE"
        to_color = bg.get("to_color") or "#FFFFFF"
        origin = bg.get("origin") or [0.3, 0.5]
        if not (isinstance(origin, (list, tuple)) and len(origin) >= 2):
            origin = [0.3, 0.5]
        ox = _safe_float(origin[0], 0.3)
        oy = _safe_float(origin[1], 0.5)
        ox = max(0.0, min(1.0, ox))
        oy = max(0.0, min(1.0, oy))

        bg_path = asset_dir / f"bg_spec_{slide_index:02d}.jpg"
        try:
            asyncio.get_event_loop()
        except RuntimeError:
            pass

        # generate_background 는 async 함수이지만 내부 작업이 동기.
        # render_design_spec 는 동기 함수이므로 직접 numpy 호출을 우회하기 위해
        # asyncio.run 대신 ppt_asset_service 의 내부 동기 헬퍼를 사용해야 하지만,
        # 안전하게 별 이벤트 루프 없이 직접 처리하기 위해 inline 으로 짧게 호출한다.
        try:
            import numpy as np  # type: ignore
            from PIL import Image as _PILImage  # noqa: F401
            from services.ppt_asset_service import (
                _hex_to_np, _make_meshgrid, _radial, _linear, _save_jpg,
            )
            W, H = 1920, 1080
            c_in = _hex_to_np(from_color)
            c_out = _hex_to_np(to_color)
            X, Y = _make_meshgrid(W, H)
            if style == "linear":
                arr = _linear(W, H, X, Y, ox, oy, 0.0, 1.2, 1.2, c_in, c_out)
            else:
                arr = _radial(W, H, X, Y, ox, oy, 0.85, 1.3, c_in, c_out)
            asset_dir.mkdir(parents=True, exist_ok=True)
            _save_jpg(arr, str(bg_path))
            slide.shapes.add_picture(
                str(bg_path), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
        except Exception as e:
            print(f"[Builder-Spec] gradient 배경 생성 실패: {e}")
        return

    # btype == "none" 또는 알 수 없음 → 흰 배경 유지


# ============ M16.A — 빌드 후 PPTX 후처리 (폰트 임베드) ============

# OOXML 네임스페이스
_PRES_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}
_CT_NS = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
_REL_NS = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}


def _resolve_font_local_path(font_meta: dict) -> Optional[str]:
    """font_meta (skill_file.fonts 의 한 항목) → 디스크 ttf/otf 절대경로.

    우선순위: local_path > url(/uploads/...) → 로컬 경로 변환.
    """
    if not isinstance(font_meta, dict):
        return None
    # 직접 local_path 가 주어진 경우
    lp = font_meta.get("local_path")
    if isinstance(lp, str) and lp and os.path.isfile(lp):
        return lp
    # url 로 시도
    url = font_meta.get("url")
    if isinstance(url, str) and url:
        if url.startswith("/uploads/"):
            rel = url[len("/uploads/"):]
            full = os.path.join(settings.UPLOAD_DIR, rel.replace("/", os.sep))
            if os.path.isfile(full):
                return full
        elif os.path.isfile(url):
            return url
    return None


def _embed_fonts_in_pptx(pptx_path: str, font_entries: list) -> int:
    """빌드 완료된 .pptx (zip) 에 폰트 파일을 임베드.

    font_entries: [{"family": str, "local_path": str}] (이미 로컬 경로 확정된 폰트만)

    OOXML 변경 사항:
      1. ppt/fonts/font<N>.ttf 추가
      2. ppt/presentation.xml 의 <p:presentation> 안에 <p:embeddedFontLst> 추가
      3. ppt/_rels/presentation.xml.rels 에 관계 4개 추가 (font/relationship)
      4. [Content_Types].xml 에 Override 추가

    반환: 실제 임베드 성공한 폰트 개수.
    """
    if not font_entries:
        return 0

    pptx_path = str(pptx_path)
    if not os.path.isfile(pptx_path):
        return 0

    # 임시 작업 디렉토리에 zip 풀기
    tmp_dir = Path(pptx_path + ".embed_tmp")
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir, ignore_errors=True)
    tmp_dir.mkdir(parents=True, exist_ok=True)

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            zf.extractall(tmp_dir)

        # 1. ppt/fonts 디렉토리에 폰트 복사
        fonts_dir = tmp_dir / "ppt" / "fonts"
        fonts_dir.mkdir(parents=True, exist_ok=True)

        # 폰트 파일명 결정 + 복사
        embedded: list[dict] = []
        for idx, fe in enumerate(font_entries, start=1):
            family = (fe.get("family") or "").strip()
            src = fe.get("local_path")
            if not family or not src or not os.path.isfile(src):
                continue
            ext = os.path.splitext(src)[1].lower()
            if ext not in (".ttf", ".otf"):
                # PPTX 임베드는 TTF/OTF 만 — 그 외 형식 스킵
                continue
            font_filename = f"font{idx}{ext}"
            target = fonts_dir / font_filename
            try:
                shutil.copy2(src, target)
            except Exception as e:
                print(f"[Embed] 폰트 복사 실패 {src}: {e}")
                continue
            embedded.append({
                "family": family,
                "filename": font_filename,
                "ext": ext,
            })

        if not embedded:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return 0

        # 2. ppt/_rels/presentation.xml.rels — 관계 추가
        rels_path = tmp_dir / "ppt" / "_rels" / "presentation.xml.rels"
        if not rels_path.is_file():
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return 0

        rels_tree = etree.parse(str(rels_path))
        rels_root = rels_tree.getroot()
        existing_ids = set()
        for r in rels_root.findall("{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
            rid = r.get("Id")
            if rid:
                existing_ids.add(rid)

        # 고유 rId 생성
        def _next_rid() -> str:
            n = 1
            while f"rId_emb_font_{n}" in existing_ids:
                n += 1
            rid = f"rId_emb_font_{n}"
            existing_ids.add(rid)
            return rid

        # OOXML font relationship type
        FONT_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
        for e in embedded:
            rid = _next_rid()
            e["rid"] = rid
            new_rel = etree.SubElement(
                rels_root,
                "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
            )
            new_rel.set("Id", rid)
            new_rel.set("Type", FONT_REL_TYPE)
            new_rel.set("Target", f"fonts/{e['filename']}")

        rels_tree.write(str(rels_path), xml_declaration=True, encoding="UTF-8", standalone=True)

        # 3. ppt/presentation.xml — embeddedFontLst 추가
        pres_path = tmp_dir / "ppt" / "presentation.xml"
        if not pres_path.is_file():
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return 0

        pres_tree = etree.parse(str(pres_path))
        pres_root = pres_tree.getroot()

        # 기존 embeddedFontLst 가 있으면 제거 후 재생성
        for child in pres_root.findall(qn("p:embeddedFontLst")):
            pres_root.remove(child)

        emb_lst = etree.SubElement(pres_root, qn("p:embeddedFontLst"))
        for e in embedded:
            emb_font = etree.SubElement(emb_lst, qn("p:embeddedFont"))
            font_el = etree.SubElement(emb_font, qn("p:font"))
            font_el.set("typeface", e["family"])
            font_el.set("panose", "020F0502020204030204")  # generic sans-serif panose
            font_el.set("pitchFamily", "34")
            font_el.set("charset", "0")
            regular_el = etree.SubElement(emb_font, qn("p:regular"))
            regular_el.set(qn("r:id"), e["rid"])

        # ppt:presentation 의 자식 순서가 까다로움 → embeddedFontLst 는
        # defaultTextStyle 뒤에 와야 함. 위에서 SubElement 로 마지막에 붙였으므로
        # 일반적으로 OK 지만, 명세상 위치 조정.
        try:
            default_text_style = pres_root.find(qn("p:defaultTextStyle"))
            if default_text_style is not None:
                # embeddedFontLst 를 defaultTextStyle 뒤로 이동
                pres_root.remove(emb_lst)
                idx = list(pres_root).index(default_text_style) + 1
                pres_root.insert(idx, emb_lst)
        except Exception:
            pass

        pres_tree.write(str(pres_path), xml_declaration=True, encoding="UTF-8", standalone=True)

        # 4. [Content_Types].xml — Default(또는 Override) 추가
        ct_path = tmp_dir / "[Content_Types].xml"
        if not ct_path.is_file():
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return 0

        ct_tree = etree.parse(str(ct_path))
        ct_root = ct_tree.getroot()
        CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
        existing_defaults = {d.get("Extension"): d.get("ContentType")
                             for d in ct_root.findall(f"{{{CT_NS}}}Default")}

        # TTF default 가 없으면 추가
        if "ttf" not in existing_defaults:
            d = etree.SubElement(ct_root, f"{{{CT_NS}}}Default")
            d.set("Extension", "ttf")
            d.set("ContentType", "application/x-fontTTF")
        if "otf" not in existing_defaults and any(e["ext"] == ".otf" for e in embedded):
            d = etree.SubElement(ct_root, f"{{{CT_NS}}}Default")
            d.set("Extension", "otf")
            d.set("ContentType", "application/vnd.ms-opentype")

        ct_tree.write(str(ct_path), xml_declaration=True, encoding="UTF-8", standalone=True)

        # 5. 새 zip 으로 패키징 (원본 덮어쓰기)
        new_pptx = str(tmp_dir) + ".out.pptx"
        with zipfile.ZipFile(new_pptx, "w", zipfile.ZIP_DEFLATED) as zf:
            for root, _dirs, files in os.walk(tmp_dir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    rel = os.path.relpath(fpath, tmp_dir).replace(os.sep, "/")
                    zf.write(fpath, rel)

        # 원본 교체
        shutil.move(new_pptx, pptx_path)
        return len(embedded)
    except Exception as e:
        print(f"[Embed] 폰트 임베드 전체 실패: {e}")
        import traceback
        traceback.print_exc()
        return 0
    finally:
        try:
            shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass


def _post_process_pptx(pptx_path: str, skill_file: dict) -> dict:
    """python-pptx 저장 후 호출하는 깊은 OOXML 후처리.

    현재 기능:
      1. skill_file.fonts 의 TTF/OTF 파일을 PPTX 에 임베드

    반환: {"embedded_fonts": int, "errors": [...] }
    """
    result = {"embedded_fonts": 0, "errors": []}
    if not isinstance(skill_file, dict):
        return result

    fonts_list = skill_file.get("fonts") or []
    if not isinstance(fonts_list, list):
        return result

    # 로컬 경로 확정된 폰트만 추출
    font_entries = []
    for f in fonts_list:
        if not isinstance(f, dict):
            continue
        family = (f.get("family") or f.get("name") or "").strip()
        if not family:
            continue
        local = _resolve_font_local_path(f)
        if not local:
            continue
        font_entries.append({"family": family, "local_path": local})

    if not font_entries:
        return result

    try:
        embedded_n = _embed_fonts_in_pptx(pptx_path, font_entries)
        result["embedded_fonts"] = embedded_n
    except Exception as e:
        result["errors"].append(f"font_embed: {e}")
        print(f"[PostProcess] 폰트 임베드 실패: {e}")

    return result


def render_design_spec(prs, design_spec: dict, skill_file: dict, asset_dir: Path):
    """단일 디자인 스펙(dict) → python-pptx 슬라이드 렌더.

    Steps:
      1. 빈 슬라이드 추가 (layout 6: blank)
      2. background 렌더 (gradient | solid | image | none)
      3. regions 순서대로 렌더 (text | shape | icon | image)
      4. 좌표 clamp + 누락 필드 안전 기본값 + 알 수 없는 region.type skip
    """
    if not isinstance(design_spec, dict):
        return
    asset_dir = Path(asset_dir)

    # 1. 빈 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 2. 배경
    slide_index = _safe_int(design_spec.get("slide_index"), 1)
    try:
        _render_background(slide, prs, design_spec, asset_dir, slide_index)
    except Exception as e:
        print(f"[Builder-Spec] 배경 렌더 실패 idx={slide_index}: {e}")

    # 3. regions
    regions = design_spec.get("regions") or []
    if not isinstance(regions, list):
        regions = []

    icon_dir = asset_dir / "_icons"

    for r in regions:
        if not isinstance(r, dict):
            continue
        rtype = (r.get("type") or "").lower()
        try:
            if rtype == "text":
                _render_text_region(slide, r, skill_file)
            elif rtype == "shape":
                _render_shape_region(slide, r)
            elif rtype == "icon":
                _render_icon_region(slide, r, icon_dir)
            elif rtype == "image":
                _render_image_region(slide, r)
            # M16.A — 새 region 타입
            elif rtype == "ghost_text":
                _render_ghost_text_region(slide, r, skill_file)
            elif rtype == "page_indicator":
                _render_page_indicator_region(slide, r, skill_file)
            elif rtype == "chip":
                _render_chip_region(slide, r, skill_file)
            elif rtype == "accent_line":
                _render_accent_line_region(slide, r)
            elif rtype == "decoration_set":
                _render_decoration_set_region(slide, r)
            else:
                print(f"[Builder-Spec] 알 수 없는 region.type: {rtype} (skip)")
        except Exception as e:
            print(f"[Builder-Spec] region 렌더 실패 type={rtype}: {e}")


async def build_pptx_styled_from_designs_stream(
    project_id: str,
    style_id: Optional[str] = None,
    design_specs: Optional[list] = None,
    skill_file: Optional[dict] = None,
    output_dir: Optional[str] = None,
):
    """design_specs 를 순서대로 렌더하며 슬라이드별 SSE 이벤트를 yield.

    Yields tuples of (event_type, payload):
      - ("slide_start", {"index", "total", "layout_hint", "title"})
      - ("slide_done",  {"index", "total"})
      - ("parsed",      {"slide_count", "object_count"})
      - ("parse_failed",{"message"})
      - ("result",      {"pptx_path", "pptx_url", "total_slides",
                         "parsed_slide_count", "parsed_object_count"})
      - ("error",       {"message"})

    design_specs / skill_file 인자가 None 이면 DB(generated_pptx_styled) 에서 로드.
    """
    db = get_db()

    # 1) design_specs 로드 (DB)
    rec = await db.generated_pptx_styled.find_one({"project_id": project_id})
    if design_specs is None:
        if not rec:
            yield ("error", {"message": "디자인이 아직 생성되지 않았습니다."})
            return
        design_specs = rec.get("design_specs") or []
    if not design_specs:
        yield ("error", {"message": "디자인이 아직 생성되지 않았습니다."})
        return

    # 2) skill_file 로드
    if skill_file is None:
        sid = style_id or (rec or {}).get("style_id")
        if not sid:
            yield ("error", {"message": "style_id 가 없습니다."})
            return
        skill_file = await ppt_style_service.generate_skill_file(sid)
        if not skill_file:
            yield ("error", {"message": "스타일 스킬 파일을 로드할 수 없습니다."})
            return

    # 3) 출력 경로
    base_dir = Path(output_dir) if output_dir else (
        Path(settings.UPLOAD_DIR) / "ppt_styled" / str(project_id)
    )
    base_dir.mkdir(parents=True, exist_ok=True)

    # 4) Presentation 객체 (16:9)
    prs = Presentation()
    prs.slide_width = Inches(CANVAS_W_IN)
    prs.slide_height = Inches(CANVAS_H_IN)

    total = len(design_specs)

    # 5) 슬라이드 루프
    for i, spec in enumerate(design_specs, start=1):
        if not isinstance(spec, dict):
            spec = {}
        layout_hint = spec.get("layout_hint", "")
        # title 추정: regions 의 가장 큰 폰트 텍스트 또는 첫 text region
        title_guess = ""
        regions = spec.get("regions") or []
        if isinstance(regions, list):
            biggest_size = -1
            for r in regions:
                if isinstance(r, dict) and r.get("type") == "text":
                    fs = _safe_int(r.get("font_size"), 0)
                    if fs > biggest_size and r.get("text"):
                        biggest_size = fs
                        title_guess = str(r.get("text", ""))[:50]

        yield ("slide_start", {
            "index": i,
            "total": total,
            "layout_hint": layout_hint,
            "title": title_guess,
        })

        try:
            render_design_spec(prs, spec, skill_file, base_dir)
        except Exception as e:
            print(f"[Builder-Spec] 슬라이드 빌드 실패 idx={i}: {e}")
            import traceback
            traceback.print_exc()
            # 빈 슬라이드 하나 추가하여 인덱스 보존
            try:
                prs.slides.add_slide(prs.slide_layouts[6])
            except Exception:
                pass

        yield ("slide_done", {"index": i, "total": total})
        await asyncio.sleep(0)

    # 6) PPTX 저장
    pptx_path = base_dir / "presentation.pptx"
    try:
        prs.save(str(pptx_path))
    except Exception as e:
        yield ("error", {"message": f"PPTX 저장 실패: {e}"})
        return

    # 6a) M16.A — 빌드 후 후처리 (폰트 임베드 등)
    post_result = {"embedded_fonts": 0, "errors": []}
    try:
        post_result = _post_process_pptx(str(pptx_path), skill_file)
        if post_result.get("embedded_fonts"):
            yield ("fonts_embedded", {"count": post_result["embedded_fonts"]})
    except Exception as e:
        print(f"[Builder-Spec] post_process 전체 실패: {e}")

    pptx_url = f"/uploads/ppt_styled/{project_id}/presentation.pptx"
    await db.generated_pptx_styled.update_one(
        {"project_id": project_id},
        {"$set": {
            "status": "built",
            "pptx_path": str(pptx_path),
            "pptx_url": pptx_url,
            "embedded_fonts_count": post_result.get("embedded_fonts", 0),
            "built_at": datetime.utcnow(),
            "updated_at": datetime.utcnow(),
        }},
    )

    # 7) M5 — 빌드된 PPTX 파싱하여 generated_slides 컬렉션에 저장
    parsed_slide_count = 0
    parsed_object_count = 0
    parse_error_msg: Optional[str] = None
    try:
        parse_result = await asyncio.to_thread(parse_pptx_to_slides, str(pptx_path))
        parsed_slides = parse_result.get("slides") or []

        await db.generated_slides.delete_many({"project_id": project_id})

        now = datetime.utcnow()
        slide_docs = []
        for idx, parsed in enumerate(parsed_slides):
            if not isinstance(parsed, dict):
                continue
            objects = parsed.get("objects") or []
            slide_doc = {
                "project_id": project_id,
                "order": idx + 1,
                "objects": objects,
                "slide_meta": parsed.get("slide_meta") or {},
                "background_image": parsed.get("background_image"),
                "items": [],
                "created_at": now,
                "updated_at": now,
            }
            slide_docs.append(slide_doc)
            parsed_object_count += len(objects)

        if slide_docs:
            await db.generated_slides.insert_many(slide_docs)
        parsed_slide_count = len(slide_docs)

        await db.generated_pptx_styled.update_one(
            {"project_id": project_id},
            {"$set": {
                "parsed_slide_count": parsed_slide_count,
                "parsed_object_count": parsed_object_count,
                "parsed_at": datetime.utcnow(),
                "updated_at": datetime.utcnow(),
            }, "$unset": {"parse_error": ""}},
        )

        yield ("parsed", {
            "slide_count": parsed_slide_count,
            "object_count": parsed_object_count,
        })
    except Exception as e:
        parse_error_msg = str(e)
        print(f"[Builder-Spec] PPTX 파싱/저장 실패: {parse_error_msg}")
        try:
            await db.generated_pptx_styled.update_one(
                {"project_id": project_id},
                {"$set": {
                    "parse_error": parse_error_msg,
                    "updated_at": datetime.utcnow(),
                }},
            )
        except Exception as e2:
            print(f"[Builder-Spec] parse_error DB 기록 실패: {e2}")
        yield ("parse_failed", {"message": parse_error_msg})

    yield ("result", {
        "pptx_path": str(pptx_path),
        "pptx_url": pptx_url,
        "total_slides": total,
        "parsed_slide_count": parsed_slide_count,
        "parsed_object_count": parsed_object_count,
    })
