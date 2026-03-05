import os
import csv
import io
from config import settings


def get_upload_path(sub_dir: str, filename: str) -> str:
    """업로드 경로 생성"""
    path = os.path.join(settings.UPLOAD_DIR, sub_dir)
    os.makedirs(path, exist_ok=True)
    return os.path.join(path, filename)


def delete_file(file_path: str) -> bool:
    """파일 삭제"""
    full_path = os.path.join(".", file_path.lstrip("/"))
    if os.path.exists(full_path):
        os.remove(full_path)
        return True
    return False


def extract_text_from_file(file_path: str, ext: str) -> str:
    """파일에서 텍스트를 추출하여 마크다운 형식으로 반환"""
    ext = ext.lower()
    try:
        if ext == ".txt":
            return _extract_txt(file_path)
        elif ext == ".csv":
            return _extract_csv(file_path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext == ".pdf":
            return _extract_pdf(file_path)
        else:
            return ""
    except Exception as e:
        return f"텍스트 추출 오류: {str(e)}"


def _extract_txt(file_path: str) -> str:
    """텍스트 파일 추출"""
    encodings = ["utf-8", "cp949", "euc-kr", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def _extract_csv(file_path: str) -> str:
    """CSV 파일을 마크다운 테이블로 변환"""
    encodings = ["utf-8", "cp949", "euc-kr", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc, newline="") as f:
                reader = csv.reader(f)
                rows = list(reader)
                if not rows:
                    return ""
                return _rows_to_markdown_table(rows)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def _extract_docx(file_path: str) -> str:
    """Word 문서에서 텍스트 추출 (마크다운 형식)"""
    from docx import Document
    doc = Document(file_path)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = (para.style.name or "").lower()
        if "heading 1" in style_name:
            parts.append(f"# {text}")
        elif "heading 2" in style_name:
            parts.append(f"## {text}")
        elif "heading 3" in style_name:
            parts.append(f"### {text}")
        elif "heading" in style_name:
            parts.append(f"#### {text}")
        elif "list" in style_name or "bullet" in style_name:
            parts.append(f"- {text}")
        else:
            parts.append(text)

    # 테이블 추출
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            parts.append("")
            parts.append(_rows_to_markdown_table(rows))

    return "\n\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    """Excel 파일을 마크다운 테이블로 변환"""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"## {sheet}")

        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c for c in cells):
                rows.append(cells)

        if rows:
            parts.append(_rows_to_markdown_table(rows))
        else:
            parts.append("(빈 시트)")

    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    """PowerPoint에서 텍스트 추출 (마크다운 형식)"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    slide_texts.append(_rows_to_markdown_table(rows))

        if slide_texts:
            parts.append(f"### 슬라이드 {i}")
            parts.append("\n".join(slide_texts))

    return "\n\n".join(parts)


def _extract_pdf(file_path: str) -> str:
    """PDF에서 텍스트 추출"""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    parts = []

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            parts.append(f"### 페이지 {i}")
            parts.append(text.strip())

    return "\n\n".join(parts)


def _rows_to_markdown_table(rows: list) -> str:
    """2D 배열을 마크다운 테이블로 변환"""
    if not rows:
        return ""

    # 첫 행을 헤더로 사용
    max_cols = max(len(row) for row in rows)

    # 모든 행의 열 수를 동일하게 맞춤
    normalized = []
    for row in rows:
        padded = list(row) + [""] * (max_cols - len(row))
        normalized.append(padded)

    lines = []
    # 헤더
    header = "| " + " | ".join(str(c).replace("|", "\\|") for c in normalized[0]) + " |"
    lines.append(header)
    # 구분선
    separator = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
    lines.append(separator)
    # 데이터 행
    for row in normalized[1:]:
        line = "| " + " | ".join(str(c).replace("|", "\\|") for c in row) + " |"
        lines.append(line)

    return "\n".join(lines)
