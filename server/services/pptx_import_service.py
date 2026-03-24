"""PPTX 파일을 파싱하여 템플릿 + 슬라이드를 MongoDB에 자동 생성하는 서비스"""

import os
import uuid
import logging
import asyncio
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from services.mongo_service import get_db
from config import settings

logger = logging.getLogger(__name__)

# ============ 캔버스 / PPTX 상수 (ppt_service.py와 동일) ============

SLIDE_W_PX = 960
SLIDE_H_PX = 540
# 기본 EMU (16:9 표준): 실제 파싱 시 PPTX 슬라이드 크기로 대체됨
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5148000

# 슬라이드 사이즈별 상수 (ppt_service.py와 동일)
SLIDE_SIZES = {
    "16:9": {"px_w": 960, "px_h": 540, "emu_w": 9144000,  "emu_h": 5148000},
    "4:3":  {"px_w": 960, "px_h": 720, "emu_w": 9144000,  "emu_h": 6858000},
    "A4":   {"px_w": 960, "px_h": 665, "emu_w": 9900000,  "emu_h": 6858000},
}

# ============ 역방향 도형 매핑 (MSO_SHAPE enum → 문자열) ============

_FORWARD_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "snip_1_rect": MSO_SHAPE.SNIP_1_RECTANGLE,
    "snip_2_diag_rect": MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
    "round_1_rect": MSO_SHAPE.ROUND_1_RECTANGLE,
    "round_2_diag_rect": MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "heptagon": MSO_SHAPE.HEPTAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "decagon": MSO_SHAPE.DECAGON,
    "dodecagon": MSO_SHAPE.DODECAGON,
    "cross": MSO_SHAPE.CROSS,
    "donut": MSO_SHAPE.DONUT,
    "no_smoking": MSO_SHAPE.NO_SYMBOL,
    "block_arc": MSO_SHAPE.BLOCK_ARC,
    "heart": MSO_SHAPE.HEART,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "sun": MSO_SHAPE.SUN,
    "moon": MSO_SHAPE.MOON,
    "cloud": MSO_SHAPE.CLOUD,
    "smiley_face": MSO_SHAPE.SMILEY_FACE,
    "folded_corner": MSO_SHAPE.FOLDED_CORNER,
    "frame": MSO_SHAPE.FRAME,
    "teardrop": MSO_SHAPE.TEAR,
    "plaque": MSO_SHAPE.PLAQUE,
    "brace_pair": MSO_SHAPE.DOUBLE_BRACE,
    "bracket_pair": MSO_SHAPE.DOUBLE_BRACKET,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "up_down_arrow": MSO_SHAPE.UP_DOWN_ARROW,
    "quad_arrow": MSO_SHAPE.QUAD_ARROW,
    "notched_right_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "home_plate": MSO_SHAPE.PENTAGON,
    "striped_right_arrow": MSO_SHAPE.STRIPED_RIGHT_ARROW,
    "bent_arrow": MSO_SHAPE.BENT_ARROW,
    "u_turn_arrow": MSO_SHAPE.U_TURN_ARROW,
    "circular_arrow": MSO_SHAPE.CIRCULAR_ARROW,
    "math_plus": MSO_SHAPE.MATH_PLUS,
    "math_minus": MSO_SHAPE.MATH_MINUS,
    "math_multiply": MSO_SHAPE.MATH_MULTIPLY,
    "math_divide": MSO_SHAPE.MATH_DIVIDE,
    "math_equal": MSO_SHAPE.MATH_EQUAL,
    "math_not_equal": MSO_SHAPE.MATH_NOT_EQUAL,
    "star_4_point": MSO_SHAPE.STAR_4_POINT,
    "star_5_point": MSO_SHAPE.STAR_5_POINT,
    "star_6_point": MSO_SHAPE.STAR_6_POINT,
    "star_8_point": MSO_SHAPE.STAR_8_POINT,
    "star_10_point": MSO_SHAPE.STAR_10_POINT,
    "star_12_point": MSO_SHAPE.STAR_12_POINT,
    "star_16_point": MSO_SHAPE.STAR_16_POINT,
    "star_24_point": MSO_SHAPE.STAR_24_POINT,
    "star_32_point": MSO_SHAPE.STAR_32_POINT,
    "explosion_1": MSO_SHAPE.EXPLOSION1,
    "explosion_2": MSO_SHAPE.EXPLOSION2,
    "wave": MSO_SHAPE.WAVE,
    "double_wave": MSO_SHAPE.DOUBLE_WAVE,
    "ribbon": MSO_SHAPE.DOWN_RIBBON,
    "wedge_rect_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "wedge_round_rect_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    "wedge_ellipse_callout": MSO_SHAPE.OVAL_CALLOUT,
    "cloud_callout": MSO_SHAPE.CLOUD_CALLOUT,
    "border_callout_1": MSO_SHAPE.LINE_CALLOUT_1,
    "border_callout_2": MSO_SHAPE.LINE_CALLOUT_2,
    "border_callout_3": MSO_SHAPE.LINE_CALLOUT_3,
}

# MSO_SHAPE enum value → 문자열 이름 (역방향 매핑)
_REVERSE_SHAPE_MAP = {v: k for k, v in _FORWARD_SHAPE_MAP.items()}

# ============ TOC / Closing 키워드 ============

_TOC_KEYWORDS = {"목차", "contents", "table of contents", "index"}
_CLOSING_KEYWORDS = {"감사", "thank", "q&a", "질의", "문의"}

# ============ 차트 타입 매핑 ============

_XL_CHART_TYPE_MAP = {}

def _map_xl_chart_type(xl_type) -> str:
    """XL_CHART_TYPE enum 값을 내부 문자열 타입으로 매핑"""
    try:
        # Bar / Column 계열
        if xl_type in (
            XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.BAR_STACKED_100,
            XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.COLUMN_STACKED_100,
            XL_CHART_TYPE.THREE_D_BAR_CLUSTERED, XL_CHART_TYPE.THREE_D_BAR_STACKED,
            XL_CHART_TYPE.THREE_D_BAR_STACKED_100,
            XL_CHART_TYPE.THREE_D_COLUMN, XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED,
            XL_CHART_TYPE.THREE_D_COLUMN_STACKED, XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100,
        ):
            return "bar"
        # Line 계열
        if xl_type in (
            XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS,
            XL_CHART_TYPE.LINE_MARKERS_STACKED, XL_CHART_TYPE.LINE_STACKED,
            XL_CHART_TYPE.LINE_MARKERS_STACKED_100, XL_CHART_TYPE.LINE_STACKED_100,
            XL_CHART_TYPE.THREE_D_LINE,
        ):
            return "line"
        # Pie 계열
        if xl_type in (
            XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED,
            XL_CHART_TYPE.THREE_D_PIE, XL_CHART_TYPE.THREE_D_PIE_EXPLODED,
            XL_CHART_TYPE.PIE_OF_PIE, XL_CHART_TYPE.BAR_OF_PIE,
        ):
            return "pie"
        # Doughnut 계열
        if xl_type in (
            XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        ):
            return "doughnut"
        # Area 계열
        if xl_type in (
            XL_CHART_TYPE.AREA, XL_CHART_TYPE.AREA_STACKED,
            XL_CHART_TYPE.AREA_STACKED_100,
            XL_CHART_TYPE.THREE_D_AREA, XL_CHART_TYPE.THREE_D_AREA_STACKED,
            XL_CHART_TYPE.THREE_D_AREA_STACKED_100,
        ):
            return "area"
        # Radar 계열
        if xl_type in (
            XL_CHART_TYPE.RADAR, XL_CHART_TYPE.RADAR_FILLED,
            XL_CHART_TYPE.RADAR_MARKERS,
        ):
            return "radar"
    except (AttributeError, TypeError, ValueError):
        pass
    return "bar"  # 기본값


# ============ 유틸리티 함수 ============


# 파싱 중 실제 슬라이드 크기를 저장하는 컨텍스트 (parse_pptx_to_slides에서 설정)
_parse_ctx = {"w_emu": SLIDE_W_EMU, "h_emu": SLIDE_H_EMU, "px_h": SLIDE_H_PX}


def _emu_to_px(emu_val, is_horizontal=True):
    """EMU 값을 캔버스 px 좌표로 변환 (ppt_service.py의 역연산)

    실제 슬라이드 크기는 _parse_ctx에서 참조 (parse_pptx_to_slides에서 설정)
    """
    if is_horizontal:
        return round(emu_val / _parse_ctx["w_emu"] * SLIDE_W_PX, 1)
    return round(emu_val / _parse_ctx["h_emu"] * _parse_ctx["px_h"], 1)


def _gen_obj_id():
    """admin.js 호환 오브젝트 ID 생성"""
    return "obj_" + uuid.uuid4().hex[:12]


def _safe_rgb(color_obj):
    """안전하게 색상 hex 문자열 추출 (테마 색상 등 예외 처리)"""
    try:
        if color_obj is not None and color_obj.rgb is not None:
            return f"#{color_obj.rgb}"
    except (AttributeError, TypeError, ValueError):
        pass
    return "#000000"


def _safe_font_color(font):
    """폰트에서 색상을 안전하게 추출 (color 속성 접근 자체가 예외를 던질 수 있음)"""
    try:
        color_obj = font.color
        return _safe_rgb(color_obj)
    except (AttributeError, TypeError, ValueError):
        return "#000000"


def _alignment_to_str(alignment):
    """PP_ALIGN enum → 문자열 변환"""
    if alignment == PP_ALIGN.CENTER:
        return "center"
    elif alignment == PP_ALIGN.RIGHT:
        return "right"
    return "left"


def _get_image_extension(content_type: str) -> str:
    """이미지 content_type → 파일 확장자 매핑"""
    type_map = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/webp": "webp",
        "image/svg+xml": "svg",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    return type_map.get(content_type, "png")


def _shape_covers_slide(left, top, width, height):
    """도형이 슬라이드 면적의 90% 이상을 차지하는지 확인 (EMU 값 기준)"""
    shape_area = width * height
    slide_area = SLIDE_W_EMU * SLIDE_H_EMU
    if slide_area == 0:
        return False
    coverage = shape_area / slide_area

    # 위치도 대략 원점 근처여야 함
    near_origin = (abs(left) < SLIDE_W_EMU * 0.05) and (abs(top) < SLIDE_H_EMU * 0.05)

    return coverage >= 0.9 and near_origin


def _save_image_blob(blob: bytes, content_type: str, subfolder: str = "images") -> str:
    """이미지 바이트를 디스크에 저장하고 URL 경로를 반환"""
    ext = _get_image_extension(content_type)
    filename = f"{uuid.uuid4().hex}.{ext}"
    save_dir = os.path.join(settings.UPLOAD_DIR, subfolder)
    os.makedirs(save_dir, exist_ok=True)
    save_path = os.path.join(save_dir, filename)

    with open(save_path, "wb") as f:
        f.write(blob)

    return f"/uploads/{subfolder}/{filename}"


# ============ 텍스트 스타일 추출 ============


def _extract_text_style(paragraph):
    """가장 긴 run(dominant)에서 텍스트 스타일 추출 (기본값 적용)"""
    style = {
        "font_family": "Arial",
        "font_size": 16,
        "color": "#000000",
        "bold": False,
        "italic": False,
        "align": _alignment_to_str(paragraph.alignment),
    }

    if not paragraph.runs:
        return style

    # 가장 텍스트가 긴 run을 dominant로 선택 (스타일 대표성)
    run = max(paragraph.runs, key=lambda r: len(r.text)) if len(paragraph.runs) > 1 else paragraph.runs[0]
    font = run.font

    try:
        if font.name:
            style["font_family"] = font.name
    except (AttributeError, TypeError):
        pass

    try:
        if font.size is not None:
            # font.size는 Emu 단위 → Pt 변환: emu / 12700
            style["font_size"] = max(1, round(font.size / 12700))
    except (AttributeError, TypeError):
        pass

    style["color"] = _safe_font_color(font)

    try:
        style["bold"] = bool(font.bold)
    except (AttributeError, TypeError):
        pass

    try:
        style["italic"] = bool(font.italic)
    except (AttributeError, TypeError):
        pass

    return style


def _extract_text_content(text_frame):
    """텍스트 프레임의 모든 paragraph를 합쳐 텍스트 반환"""
    lines = []
    for para in text_frame.paragraphs:
        lines.append(para.text)
    return "\n".join(lines)


def _get_largest_font_size(text_frame):
    """텍스트 프레임에서 가장 큰 폰트 사이즈 (Pt)를 반환"""
    max_size = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.size is not None:
                    pt_size = run.font.size / 12700
                    if pt_size > max_size:
                        max_size = pt_size
            except (AttributeError, TypeError):
                pass
    return max_size


# ============ 개별 도형 추출 ============


def _extract_text_shape(shape, z_index: int) -> dict:
    """텍스트 도형 오브젝트 추출"""
    tf = shape.text_frame
    text_content = _extract_text_content(tf)

    # 가장 텍스트가 많은 paragraph에서 스타일 추출 (대표 스타일)
    best_para = None
    best_len = 0
    for para in tf.paragraphs:
        para_len = len(para.text)
        if para_len > best_len:
            best_len = para_len
            best_para = para
    if best_para is None:
        best_para = tf.paragraphs[0] if tf.paragraphs else None
    text_style = _extract_text_style(best_para) if best_para else {
        "font_family": "Arial",
        "font_size": 16,
        "color": "#000000",
        "bold": False,
        "italic": False,
        "align": "left",
    }

    return {
        "obj_id": _gen_obj_id(),
        "obj_type": "text",
        "x": _emu_to_px(shape.left, is_horizontal=True),
        "y": _emu_to_px(shape.top, is_horizontal=False),
        "width": _emu_to_px(shape.width, is_horizontal=True),
        "height": _emu_to_px(shape.height, is_horizontal=False),
        "text_content": text_content,
        "text_style": text_style,
        "z_index": z_index,
        "role": None,       # 나중에 assign_roles에서 설정
        "placeholder": None,  # 나중에 assign_roles에서 설정
    }


def _extract_image_shape(shape, z_index: int) -> dict | None:
    """이미지 도형 오브젝트 추출 (이미지를 디스크에 저장)"""
    try:
        image = shape.image
        blob = image.blob
        content_type = image.content_type
        image_url = _save_image_blob(blob, content_type, subfolder="images")

        return {
            "obj_id": _gen_obj_id(),
            "obj_type": "image",
            "x": _emu_to_px(shape.left, is_horizontal=True),
            "y": _emu_to_px(shape.top, is_horizontal=False),
            "width": _emu_to_px(shape.width, is_horizontal=True),
            "height": _emu_to_px(shape.height, is_horizontal=False),
            "image_url": image_url,
            "image_fit": "contain",
            "z_index": z_index,
            "role": None,
            "placeholder": None,
        }
    except Exception as e:
        logger.warning(f"이미지 도형 추출 실패: {e}")
        return None


def _extract_auto_shape(shape, z_index: int) -> dict | None:
    """자동 도형 (AutoShape) 오브젝트 추출"""
    try:
        # MSO_SHAPE enum → 문자열 이름 변환
        shape_type_str = "rectangle"
        try:
            auto_shape_type = shape.auto_shape_type
            shape_type_str = _REVERSE_SHAPE_MAP.get(auto_shape_type, "rectangle")
        except (AttributeError, TypeError, ValueError):
            pass

        # 채우기 색상 / 투명도 추출
        fill_color = "#4A90D9"
        fill_opacity = 1.0
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    fore_color = fill.fore_color
                    fill_color = _safe_rgb(fore_color)
                except (AttributeError, TypeError, ValueError):
                    # 그라데이션 채우기 → 첫 번째 그라데이션 스톱 색상 사용
                    try:
                        if hasattr(fill, 'gradient_stops') and fill.gradient_stops:
                            first_stop = fill.gradient_stops[0]
                            fill_color = _safe_rgb(first_stop.color)
                    except (AttributeError, TypeError, ValueError, IndexError):
                        pass
            else:
                # fill.type이 None → 투명 채우기
                fill_opacity = 0.0
        except (AttributeError, TypeError, ValueError):
            pass

        # 테두리 색상 / 두께 추출
        stroke_color = "#333333"
        stroke_width = 2
        try:
            line = shape.line
            if line.color and line.color.rgb:
                stroke_color = f"#{line.color.rgb}"
        except (AttributeError, TypeError, ValueError):
            pass
        try:
            if shape.line.width is not None:
                stroke_width = max(0, round(shape.line.width / 12700))  # EMU → Pt
        except (AttributeError, TypeError, ValueError):
            pass

        obj = {
            "obj_id": _gen_obj_id(),
            "obj_type": "shape",
            "x": _emu_to_px(shape.left, is_horizontal=True),
            "y": _emu_to_px(shape.top, is_horizontal=False),
            "width": _emu_to_px(shape.width, is_horizontal=True),
            "height": _emu_to_px(shape.height, is_horizontal=False),
            "shape_style": {
                "shape_type": shape_type_str,
                "fill_color": fill_color,
                "fill_opacity": fill_opacity,
                "stroke_color": stroke_color,
                "stroke_width": stroke_width,
                "stroke_dash": "solid",
            },
            "z_index": z_index,
            "role": None,
            "placeholder": None,
        }

        # 텍스트가 포함된 도형 → 텍스트 오브젝트로 변환
        if shape.has_text_frame:
            text_content = _extract_text_content(shape.text_frame)
            if text_content.strip():
                first_para = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else None
                text_style = _extract_text_style(first_para) if first_para else {
                    "font_family": "Arial", "font_size": 16, "color": "#000000",
                    "bold": False, "italic": False, "align": "left",
                }
                return {
                    "obj_id": obj["obj_id"],
                    "obj_type": "text",
                    "x": obj["x"],
                    "y": obj["y"],
                    "width": obj["width"],
                    "height": obj["height"],
                    "text_content": text_content,
                    "text_style": text_style,
                    "z_index": z_index,
                    "role": None,
                    "placeholder": None,
                }

        return obj
    except Exception as e:
        logger.warning(f"자동 도형 추출 실패: {e}")
        return None


def _extract_table_as_text_fallback(shape, z_index: int) -> dict | None:
    """테이블을 텍스트 오브젝트로 변환 (내용만 추출) - 폴백용"""
    try:
        table = shape.table
        rows_text = []
        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                cells_text.append(cell.text.strip())
            rows_text.append(" | ".join(cells_text))

        text_content = "\n".join(rows_text)
        if not text_content.strip():
            return None

        return {
            "obj_id": _gen_obj_id(),
            "obj_type": "text",
            "x": _emu_to_px(shape.left, is_horizontal=True),
            "y": _emu_to_px(shape.top, is_horizontal=False),
            "width": _emu_to_px(shape.width, is_horizontal=True),
            "height": _emu_to_px(shape.height, is_horizontal=False),
            "text_content": text_content,
            "text_style": {
                "font_family": "Arial",
                "font_size": 12,
                "color": "#000000",
                "bold": False,
                "italic": False,
                "align": "left",
            },
            "z_index": z_index,
            "role": "description",
            "placeholder": None,
        }
    except Exception as e:
        logger.warning(f"테이블 추출 실패: {e}")
        return None


# ============ 차트 도형 추출 ============


def _extract_chart_shape(shape, z_index: int) -> dict | None:
    """차트 도형 오브젝트 추출"""
    try:
        chart = shape.chart
        chart_type_str = "bar"
        try:
            plot = chart.plots[0]
            chart_type_str = _map_xl_chart_type(plot.chart_type)
        except (IndexError, AttributeError, TypeError):
            pass

        # 라벨 추출
        labels = []
        try:
            plot = chart.plots[0]
            cats = plot.categories
            if cats is not None:
                labels = list(cats.flattened_labels) if hasattr(cats, 'flattened_labels') else []
        except (IndexError, AttributeError, TypeError):
            pass

        # 데이터셋 추출
        datasets = []
        color_scheme = []
        try:
            for series in chart.series:
                ds = {
                    "label": str(series.name) if series.name else "",
                    "data": [],
                }
                try:
                    ds["data"] = [v for v in series.values]
                except (AttributeError, TypeError):
                    pass

                # 시리즈 색상 추출
                try:
                    fill = series.format.fill
                    if fill.type is not None:
                        c = _safe_rgb(fill.fore_color)
                        color_scheme.append(c)
                except (AttributeError, TypeError, ValueError):
                    pass

                datasets.append(ds)
        except (AttributeError, TypeError):
            pass

        # 차트 제목 추출
        chart_title = ""
        try:
            if chart.has_title and chart.chart_title is not None:
                chart_title = chart.chart_title.text_frame.text if chart.chart_title.has_text_frame else ""
        except (AttributeError, TypeError):
            pass

        # show_legend
        show_legend = True
        try:
            show_legend = bool(chart.has_legend)
        except (AttributeError, TypeError):
            pass

        if not color_scheme:
            color_scheme = ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000", "#5B9BD5", "#70AD47"]

        return {
            "obj_id": _gen_obj_id(),
            "obj_type": "chart",
            "x": _emu_to_px(shape.left, is_horizontal=True),
            "y": _emu_to_px(shape.top, is_horizontal=False),
            "width": _emu_to_px(shape.width, is_horizontal=True),
            "height": _emu_to_px(shape.height, is_horizontal=False),
            "chart_style": {
                "chart_type": chart_type_str,
                "chart_data": {
                    "labels": labels,
                    "datasets": datasets,
                },
                "title": chart_title,
                "show_legend": show_legend,
                "show_grid": True,
                "color_scheme": color_scheme,
            },
            "z_index": z_index,
            "role": "chart",
            "placeholder": None,
        }
    except Exception as e:
        logger.warning(f"차트 도형 추출 실패: {e}")
        return None


# ============ 테이블 도형 추출 ============


def _extract_table_shape(shape, z_index: int) -> dict | None:
    """테이블을 table 오브젝트로 추출 (구조화된 데이터 포함)"""
    try:
        table = shape.table
        rows_count = len(table.rows)
        cols_count = len(table.columns)

        # 2D 데이터 배열 추출
        data = []
        cell_styles = {}
        merged_cells = []

        for r in range(rows_count):
            row_data = []
            for c in range(cols_count):
                cell = table.cell(r, c)
                row_data.append(cell.text.strip())

                # 셀 스타일 추출
                cs = {}
                try:
                    cell_fill = cell.fill
                    if cell_fill.type is not None:
                        cs["bg_color"] = _safe_rgb(cell_fill.fore_color)
                except (AttributeError, TypeError, ValueError):
                    pass

                # 텍스트 색상 / 볼드
                try:
                    if cell.text_frame and cell.text_frame.paragraphs:
                        first_para = cell.text_frame.paragraphs[0]
                        if first_para.runs:
                            run = first_para.runs[0]
                            cs["text_color"] = _safe_font_color(run.font)
                            try:
                                cs["bold"] = bool(run.font.bold)
                            except (AttributeError, TypeError):
                                pass
                except (AttributeError, TypeError):
                    pass

                if cs:
                    cell_styles[f"{r}_{c}"] = cs

                # 병합 셀 감지
                try:
                    if cell.is_merge_origin:
                        merged_cells.append({
                            "start_row": r,
                            "start_col": c,
                            "end_row": r + cell.span_height - 1,
                            "end_col": c + cell.span_width - 1,
                        })
                except (AttributeError, TypeError):
                    pass

            data.append(row_data)

        if not data:
            return _extract_table_as_text_fallback(shape, z_index)

        # 헤더 감지 (첫 번째 행의 볼드 또는 배경색)
        header_row = False
        header_bg_color = "#4472C4"
        header_text_color = "#FFFFFF"
        if rows_count > 1:
            # 첫 번째 행의 셀 스타일 확인
            first_row_styles = [cell_styles.get(f"0_{c}", {}) for c in range(cols_count)]
            has_bg = any(s.get("bg_color") for s in first_row_styles)
            has_bold = any(s.get("bold") for s in first_row_styles)
            if has_bg or has_bold:
                header_row = True
                for s in first_row_styles:
                    if s.get("bg_color"):
                        header_bg_color = s["bg_color"]
                        break
                for s in first_row_styles:
                    if s.get("text_color"):
                        header_text_color = s["text_color"]
                        break

        # 테두리 색상 추출 (XML에서)
        border_color = "#BFBFBF"
        try:
            tbl_elem = shape._element
            tc_elems = tbl_elem.findall(f".//{qn('a:tcPr')}")
            if tc_elems:
                for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                    border = tc_elems[0].find(f"{qn(border_tag)}")
                    if border is not None:
                        srgb = border.find(f".//{qn('a:srgbClr')}")
                        if srgb is not None:
                            val = srgb.get("val")
                            if val:
                                border_color = f"#{val}"
                                break
        except (AttributeError, TypeError):
            pass

        # 폰트 정보 추출 (샘플 셀에서)
        font_family = "Arial"
        font_size = 11
        try:
            for r in range(min(rows_count, 3)):
                for c in range(min(cols_count, 3)):
                    cell = table.cell(r, c)
                    if cell.text_frame and cell.text_frame.paragraphs:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    font_family = run.font.name
                                if run.font.size is not None:
                                    font_size = max(1, round(run.font.size / 12700))
                                break
                            if font_family != "Arial":
                                break
                    if font_family != "Arial":
                        break
                if font_family != "Arial":
                    break
        except (AttributeError, TypeError):
            pass

        return {
            "obj_id": _gen_obj_id(),
            "obj_type": "table",
            "x": _emu_to_px(shape.left, is_horizontal=True),
            "y": _emu_to_px(shape.top, is_horizontal=False),
            "width": _emu_to_px(shape.width, is_horizontal=True),
            "height": _emu_to_px(shape.height, is_horizontal=False),
            "table_style": {
                "rows": rows_count,
                "cols": cols_count,
                "data": data,
                "header_row": header_row,
                "banded_rows": False,
                "banded_cols": False,
                "cell_styles": cell_styles,
                "border_color": border_color,
                "border_width": 1,
                "header_bg_color": header_bg_color,
                "header_text_color": header_text_color,
                "font_family": font_family,
                "font_size": font_size,
                "merged_cells": merged_cells,
            },
            "z_index": z_index,
            "role": "table",
            "placeholder": None,
        }
    except Exception as e:
        logger.warning(f"테이블 도형 구조화 추출 실패, 폴백: {e}")
        return _extract_table_as_text_fallback(shape, z_index)


# ============ 그룹 도형 추출 ============


def _extract_group_shape(shape, z_index: int) -> list[dict]:
    """그룹 도형 내부의 자식 도형들을 재귀적으로 추출

    그룹 좌표를 절대 좌표로 변환하여 반환합니다.
    Returns: list[dict] - 추출된 오브젝트 리스트
    """
    results = []
    try:
        group_left = shape.left
        group_top = shape.top
        group_width = shape.width
        group_height = shape.height

        # 그룹 내부 좌표계 (chOff / chExt) 추출
        grpSpPr = shape._element.find(qn('p:grpSpPr'))
        if grpSpPr is None:
            grpSpPr = shape._element

        xfrm = grpSpPr.find(qn('a:xfrm'))
        if xfrm is None:
            # grpSp 자체의 xfrm
            xfrm = shape._element.find(f".//{qn('a:xfrm')}")

        ch_off_x = 0
        ch_off_y = 0
        ch_ext_cx = group_width
        ch_ext_cy = group_height

        if xfrm is not None:
            ch_off = xfrm.find(qn('a:chOff'))
            ch_ext = xfrm.find(qn('a:chExt'))
            if ch_off is not None:
                ch_off_x = int(ch_off.get('x', '0'))
                ch_off_y = int(ch_off.get('y', '0'))
            if ch_ext is not None:
                ch_ext_cx = int(ch_ext.get('cx', str(group_width)))
                ch_ext_cy = int(ch_ext.get('cy', str(group_height)))

        # 스케일 팩터 계산 (0 나누기 방지)
        scale_x = group_width / ch_ext_cx if ch_ext_cx > 0 else 1.0
        scale_y = group_height / ch_ext_cy if ch_ext_cy > 0 else 1.0

        child_z = z_index
        for child_shape in shape.shapes:
            try:
                # 재귀 처리: 자식이 그룹이면 재귀 호출
                is_group = False
                try:
                    if child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        is_group = True
                except (AttributeError, TypeError):
                    pass

                if is_group:
                    # 그룹 내부 좌표 → 절대 좌표 변환 후 재귀
                    child_results = _extract_group_shape(child_shape, child_z)
                    # 좌표 변환 (자식 그룹의 좌표는 자식 좌표계 기준)
                    for obj in child_results:
                        # 자식 결과의 좌표는 이미 px 단위이므로 EMU로 역변환 → 그룹 변환 → px 변환이 필요
                        # 하지만 자식 그룹이 재귀적으로 처리되므로 그대로 추가
                        results.append(obj)
                        child_z += 1
                    continue

                result = _extract_shape(child_shape, child_z)
                if result is None:
                    continue

                # 결과가 리스트인 경우 (중첩 그룹)
                objs = result if isinstance(result, list) else [result]
                for obj in objs:
                    # 자식 도형의 좌표를 절대 좌표로 변환
                    # child 좌표는 _extract_shape에서 이미 px로 변환됨
                    # 원래 EMU 좌표를 다시 가져와서 변환해야 함
                    try:
                        child_left = child_shape.left
                        child_top = child_shape.top
                        child_w = child_shape.width
                        child_h = child_shape.height

                        abs_x = group_left + int((child_left - ch_off_x) * scale_x)
                        abs_y = group_top + int((child_top - ch_off_y) * scale_y)
                        abs_w = int(child_w * scale_x)
                        abs_h = int(child_h * scale_y)

                        obj["x"] = _emu_to_px(abs_x, is_horizontal=True)
                        obj["y"] = _emu_to_px(abs_y, is_horizontal=False)
                        obj["width"] = _emu_to_px(abs_w, is_horizontal=True)
                        obj["height"] = _emu_to_px(abs_h, is_horizontal=False)
                    except (AttributeError, TypeError):
                        pass

                    results.append(obj)
                    child_z += 1

            except Exception as e:
                logger.warning(f"그룹 자식 도형 추출 실패: {e}")
                continue

    except Exception as e:
        logger.warning(f"그룹 도형 추출 실패: {e}")

    return results


# ============ 배경 이미지 추출 ============


def _extract_slide_background(slide) -> str | None:
    """슬라이드 배경 이미지를 추출하여 저장, URL 반환"""
    # 방법 1: slide.background.fill에서 이미지 추출
    try:
        bg = slide.background
        fill = bg.fill
        # fill type이 picture인 경우
        if fill.type is not None:
            try:
                # XML에서 직접 blipFill 확인
                bg_elem = bg._element
                blip_fills = bg_elem.findall(f".//{qn('a:blip')}")
                if blip_fills:
                    for blip in blip_fills:
                        r_embed = blip.get(qn("r:embed"))
                        if r_embed:
                            part = slide.part.related_parts.get(r_embed)
                            if part:
                                blob = part.blob
                                content_type = part.content_type or "image/png"
                                return _save_image_blob(blob, content_type, subfolder="backgrounds")
            except Exception:
                pass
    except Exception:
        pass

    return None


def _check_first_shape_as_background(shapes, total_shapes: int) -> tuple[str | None, bool]:
    """첫 번째 도형이 슬라이드 전체를 덮는 이미지인지 확인
    Returns: (bg_url 또는 None, 배경으로 사용됨 여부)
    """
    if total_shapes == 0:
        return None, False

    first_shape = shapes[0]

    # 이미지인지 확인
    try:
        shape_type = first_shape.shape_type
        is_picture = (shape_type == MSO_SHAPE_TYPE.PICTURE)
    except (AttributeError, TypeError):
        is_picture = False

    if not is_picture:
        try:
            _ = first_shape.image
            is_picture = True
        except Exception:
            is_picture = False

    if not is_picture:
        return None, False

    # 슬라이드 면적의 90% 이상을 차지하는지 확인
    try:
        if _shape_covers_slide(first_shape.left, first_shape.top,
                                first_shape.width, first_shape.height):
            image = first_shape.image
            blob = image.blob
            content_type = image.content_type
            bg_url = _save_image_blob(blob, content_type, subfolder="backgrounds")
            return bg_url, True
    except Exception:
        pass

    return None, False


# ============ 역할 할당 (Role Assignment) ============


def _assign_roles(objects: list[dict], content_type: str = "body"):
    """텍스트 오브젝트에 role과 placeholder를 자동 할당

    슬라이드 유형별 역할 할당 전략:
    - cover/title_slide : 가장 큰 폰트 → 제목, 그 다음 → 부제목, 나머지 → 내용
    - toc              : 가장 큰 폰트 → 제목, 나머지 → 내용 (목차 항목)
    - section_divider  : 가장 큰 폰트 → 제목, 나머지 → 부제목
    - body             : 가장 큰 폰트 → 제목, 두 번째 큰 폰트 → 부제목, 나머지 → 내용
    - closing          : 가장 큰 폰트 → 제목, 나머지 → 부제목

    공통 규칙:
    - Y 위치 기준으로 공간 순서를 고려
    - 폰트 크기로 중요도 판단 (큰 폰트 = 제목급)
    - 상단의 아주 작은 텍스트(≤10pt, 상위 10%) → governance
    """
    text_objs = [o for o in objects if o["obj_type"] == "text"]

    if not text_objs:
        return

    # 1단계: Y 위치 순서로 정렬 (위→아래)
    text_by_pos = sorted(text_objs, key=lambda o: o.get("y", 0))

    # 2단계: 가장 큰 폰트를 가진 오브젝트 찾기 (제목 후보)
    largest_obj = max(text_objs, key=lambda o: o.get("text_style", {}).get("font_size", 0))
    largest_font = largest_obj.get("text_style", {}).get("font_size", 16)

    # 3단계: 역할 할당
    title_obj = None
    subtitle_count = 0
    desc_count = 0

    # -- 제목 결정: 가장 큰 폰트 오브젝트 --
    if largest_font >= 14:
        title_obj = largest_obj
        title_obj["role"] = "title"
        title_obj["placeholder"] = "title_1"

    # -- 나머지 오브젝트에 역할 할당 (Y 위치 순서대로) --
    for obj in text_by_pos:
        if obj is title_obj:
            continue
        if obj.get("role") is not None:
            continue

        font_size = obj.get("text_style", {}).get("font_size", 16)
        y_pos = obj.get("y", 0)

        # governance: 상단 10% 이내, 매우 작은 폰트(≤10pt)
        if font_size <= 10 and y_pos < SLIDE_H_PX * 0.10:
            obj["role"] = "governance"
            obj["placeholder"] = "governance_1"
            continue

        # 슬라이드 유형별 부제목 vs 내용 분류
        if content_type in ("title_slide", "cover"):
            # 커버: 제목 외 → 부제목 (최대 2개), 나머지 → 내용
            if subtitle_count < 2:
                subtitle_count += 1
                obj["role"] = "subtitle"
                obj["placeholder"] = f"subtitle_{subtitle_count}"
            else:
                desc_count += 1
                obj["role"] = "description"
                obj["placeholder"] = f"desc_{desc_count}"

        elif content_type == "section_divider":
            # 간지: 제목 외 모두 → 부제목
            subtitle_count += 1
            obj["role"] = "subtitle"
            obj["placeholder"] = f"subtitle_{subtitle_count}"

        elif content_type == "closing":
            # 마무리: 제목 외 → 부제목 (최대 2개), 나머지 → 내용
            if subtitle_count < 2:
                subtitle_count += 1
                obj["role"] = "subtitle"
                obj["placeholder"] = f"subtitle_{subtitle_count}"
            else:
                desc_count += 1
                obj["role"] = "description"
                obj["placeholder"] = f"desc_{desc_count}"

        elif content_type == "toc":
            # 목차: 제목 외 모두 → 내용 (목차 항목)
            desc_count += 1
            obj["role"] = "description"
            obj["placeholder"] = f"desc_{desc_count}"

        else:
            # body (본문): 제목 다음 큰 폰트 1개 → 부제목, 나머지 → 내용
            if subtitle_count == 0 and font_size >= largest_font * 0.6 and font_size >= 14:
                subtitle_count += 1
                obj["role"] = "subtitle"
                obj["placeholder"] = f"subtitle_{subtitle_count}"
            else:
                desc_count += 1
                obj["role"] = "description"
                obj["placeholder"] = f"desc_{desc_count}"

    # 제목이 미할당인 경우 (모든 폰트가 14pt 미만)
    if title_obj is None and text_by_pos:
        # 맨 위 텍스트를 제목으로 강제 지정
        first = text_by_pos[0]
        if first.get("role") == "description":
            desc_count -= 1
        elif first.get("role") == "subtitle":
            subtitle_count -= 1
        first["role"] = "title"
        first["placeholder"] = "title_1"

    # role이 아직 None인 텍스트 오브젝트 처리
    for obj in objects:
        if obj["obj_type"] == "text" and obj.get("role") is None:
            desc_count += 1
            obj["role"] = "description"
            obj["placeholder"] = f"desc_{desc_count}"


# ============ 슬라이드 분류 ============


def _classify_slide(slide_index: int, total_slides: int, objects: list[dict]) -> dict:
    """슬라이드를 content_type과 layout으로 분류

    Returns: {"content_type": str, "layout": str}
    """
    # 모든 텍스트 합치기 (분류 키워드 탐지용)
    all_text = ""
    text_objs = []
    max_font_size = 0

    for obj in objects:
        if obj["obj_type"] == "text":
            text_objs.append(obj)
            text_content = obj.get("text_content", "")
            all_text += " " + text_content
            font_size = obj.get("text_style", {}).get("font_size", 16)
            if font_size > max_font_size:
                max_font_size = font_size

    all_text_lower = all_text.lower().strip()

    # 1. 첫 번째 슬라이드 → title_slide / cover
    if slide_index == 0:
        return {"content_type": "title_slide", "layout": "cover"}

    # 2. 마지막 슬라이드 (총 3개 이상) → closing
    if slide_index == total_slides - 1 and total_slides > 2:
        return {"content_type": "closing", "layout": "closing"}

    # 3. 목차 키워드 확인
    for kw in _TOC_KEYWORDS:
        if kw in all_text_lower:
            return {"content_type": "toc", "layout": "numbered_list"}

    # 4. 마감 키워드 확인 (감사, Thank, Q&A 등)
    for kw in _CLOSING_KEYWORDS:
        if kw in all_text_lower:
            return {"content_type": "closing", "layout": "closing"}

    # 5. 텍스트 박스 적고 (1-2개) 큰 폰트 (>24pt) → section_divider
    #    단, 테이블이나 차트가 있는 슬라이드는 section_divider가 아닌 body로 분류
    has_data_objects = any(
        obj.get("obj_type") in ("table", "chart") for obj in objects
    )
    if 1 <= len(text_objs) <= 2 and max_font_size > 24 and not has_data_objects:
        return {"content_type": "section_divider", "layout": "divider"}

    # 6. 기본값 → body / single_column
    return {"content_type": "body", "layout": "single_column"}


def _build_slide_meta(classification: dict, objects: list[dict]) -> dict:
    """슬라이드 메타 정보 생성"""
    has_title = False
    has_governance = False
    description_count = 0
    table_count = 0
    chart_count = 0

    for obj in objects:
        role = obj.get("role")
        obj_type = obj.get("obj_type")
        if role == "title":
            has_title = True
        elif role == "governance":
            has_governance = True
        elif role == "description":
            description_count += 1
        if obj_type == "table":
            table_count += 1
        elif obj_type == "chart":
            chart_count += 1

    return {
        "content_type": classification["content_type"],
        "layout": classification["layout"],
        "has_title": has_title,
        "has_governance": has_governance,
        "description_count": description_count,
        "table_count": table_count,
        "chart_count": chart_count,
    }


# ============ 슬라이드 파싱 ============


def _parse_slide(pptx_slide, slide_index: int, total_slides: int) -> dict:
    """하나의 PPTX 슬라이드를 파싱하여 오브젝트 리스트 + 분류 정보 반환"""
    objects = []
    z_index = 10

    # 배경 이미지 추출
    bg_url = _extract_slide_background(pptx_slide)

    # 도형 목록 가져오기
    shapes_list = list(pptx_slide.shapes)

    # 첫 번째 도형이 전체 슬라이드를 덮는 배경 이미지인지 확인
    first_shape_bg, is_first_bg = _check_first_shape_as_background(shapes_list, len(shapes_list))
    if first_shape_bg and not bg_url:
        bg_url = first_shape_bg

    # 도형 순회 (첫 번째가 배경이면 건너뜀)
    start_idx = 1 if is_first_bg else 0

    for shape in shapes_list[start_idx:]:
        try:
            result = _extract_shape(shape, z_index)
            if result is not None:
                if isinstance(result, list):
                    for obj in result:
                        objects.append(obj)
                        z_index += 1
                else:
                    objects.append(result)
                    z_index += 1
        except Exception as e:
            logger.warning(f"슬라이드 {slide_index + 1} 도형 추출 실패: {e}")
            continue

    # 슬라이드 분류 (역할 할당보다 먼저 수행)
    classification = _classify_slide(slide_index, total_slides, objects)

    # 텍스트 오브젝트에 role 할당 (분류 결과 활용)
    _assign_roles(objects, content_type=classification["content_type"])

    # 슬라이드 메타 정보
    slide_meta = _build_slide_meta(classification, objects)

    # 텍스트 오브젝트의 실제 내용을 플레이스홀더로 교체 (분류/역할 할당 이후)
    for obj in objects:
        if obj["obj_type"] == "text":
            obj["text_content"] = "텍스트를 입력하세요"
            # 제목은 최소 높이 60px, 부제목은 최소 높이 50px 보장
            role = obj.get("role")
            if role == "title" and obj.get("height", 0) < 60:
                obj["height"] = 60
            elif role == "subtitle" and obj.get("height", 0) < 50:
                obj["height"] = 50
            elif role == "description" and obj.get("height", 0) < 45:
                obj["height"] = 45

    return {
        "objects": objects,
        "classification": classification,
        "slide_meta": slide_meta,
        "background_image": bg_url,
    }


def _extract_shape(shape, z_index: int) -> dict | list | None:
    """도형 타입을 판별하여 적절한 추출 함수 호출

    Returns:
        dict - 단일 오브젝트
        list[dict] - 그룹 도형에서 추출된 여러 오브젝트
        None - 건너뛰기
    """
    # 그룹 도형 → 재귀적으로 추출하여 리스트 반환
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_objects = _extract_group_shape(shape, z_index)
            return group_objects if group_objects else None
    except (AttributeError, TypeError):
        pass

    # 이미지 도형 확인 (MSO_SHAPE_TYPE.PICTURE 또는 .image 속성)
    is_picture = False
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            is_picture = True
    except (AttributeError, TypeError):
        pass
    if not is_picture:
        try:
            _ = shape.image
            is_picture = True
        except Exception:
            pass

    if is_picture:
        return _extract_image_shape(shape, z_index)

    # 차트 (테이블보다 먼저 확인)
    try:
        if shape.has_chart:
            return _extract_chart_shape(shape, z_index)
    except (AttributeError, TypeError):
        pass

    # 테이블 (구조화된 추출)
    try:
        if shape.has_table:
            return _extract_table_shape(shape, z_index)
    except (AttributeError, TypeError):
        pass

    # 텍스트 프레임이 있는 도형 → 텍스트 오브젝트로 우선 처리
    has_text = False
    try:
        if shape.has_text_frame:
            text_content = _extract_text_content(shape.text_frame)
            if text_content.strip():
                has_text = True
    except (AttributeError, TypeError):
        pass
    if not has_text:
        try:
            if hasattr(shape, "text_frame"):
                text_content = _extract_text_content(shape.text_frame)
                if text_content.strip():
                    has_text = True
        except (AttributeError, TypeError):
            pass

    if has_text:
        # AutoShape에 텍스트 포함 → _extract_auto_shape가 텍스트로 변환
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return _extract_auto_shape(shape, z_index)
        except (AttributeError, TypeError):
            pass
        # 일반 텍스트 프레임
        return _extract_text_shape(shape, z_index)

    # 텍스트 없는 자동 도형 → shape 그대로 유지
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return _extract_auto_shape(shape, z_index)
    except (AttributeError, TypeError):
        pass

    # 기타 도형 → None (건너뛰기)
    return None


# ============ 메인 함수 ============


def parse_pptx_to_slides(file_path: str) -> dict:
    """PPTX 파일을 파싱하여 슬라이드 데이터를 반환 (DB 저장 없음)

    Args:
        file_path: PPTX 파일 경로

    Returns:
        {
            "slides": [{"objects": [...], "classification": {...}, "slide_meta": {...}, "background_image": str|None}],
            "slide_size": "16:9",
            "background_image": str|None,
            "total_slides": int
        }
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {file_path}")

    prs = Presentation(file_path)
    slides = list(prs.slides)
    total_slides = len(slides)

    if total_slides == 0:
        raise ValueError("슬라이드가 없는 PPTX 파일입니다")

    # 슬라이드 사이즈 감지
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    ratio = slide_width / slide_height if slide_height > 0 else 1.78
    if ratio > 1.6:
        slide_size = "16:9"
    elif ratio > 1.2:
        slide_size = "4:3"
    else:
        slide_size = "A4"

    # 실제 슬라이드 크기로 EMU 변환 컨텍스트 설정
    sz = SLIDE_SIZES.get(slide_size, SLIDE_SIZES["16:9"])
    _parse_ctx["w_emu"] = int(slide_width) if slide_width else sz["emu_w"]
    _parse_ctx["h_emu"] = int(slide_height) if slide_height else sz["emu_h"]
    _parse_ctx["px_h"] = sz["px_h"]

    # 슬라이드 파싱
    parsed_slides = []
    for idx, pptx_slide in enumerate(slides):
        try:
            parsed = _parse_slide(pptx_slide, idx, total_slides)
            parsed_slides.append(parsed)
        except Exception as e:
            logger.error(f"슬라이드 {idx + 1} 파싱 실패: {e}")
            parsed_slides.append({
                "objects": [],
                "classification": {"content_type": "body", "layout": "single_column"},
                "slide_meta": {
                    "content_type": "body",
                    "layout": "single_column",
                    "has_title": False,
                    "has_governance": False,
                    "description_count": 0,
                    "table_count": 0,
                    "chart_count": 0,
                },
                "background_image": None,
            })

    template_bg = parsed_slides[0]["background_image"] if parsed_slides else None

    return {
        "slides": parsed_slides,
        "slide_size": slide_size,
        "background_image": template_bg,
        "total_slides": total_slides,
    }


async def import_pptx_as_template(file_path: str, template_name: str, user_key: str) -> dict:
    """PPTX 파일을 분석하여 템플릿 + 슬라이드를 자동 생성

    Args:
        file_path: PPTX 파일 경로
        template_name: 생성할 템플릿 이름
        user_key: 생성자 사용자 키

    Returns:
        {"template_id": str, "slides_count": int, "classification": [...]}
    """
    # PPTX 파싱 (DB 저장 없음) - 동기 함수를 스레드풀에서 실행
    parse_result = await asyncio.to_thread(parse_pptx_to_slides, file_path)
    parsed_slides = parse_result["slides"]
    total_slides = parse_result["total_slides"]
    template_bg = parse_result["background_image"]

    db = get_db()
    now = datetime.utcnow()

    # 분류 정보 수집
    classifications = []
    for idx, parsed in enumerate(parsed_slides):
        classifications.append({
            "slide_order": idx + 1,
            "content_type": parsed["classification"]["content_type"],
            "layout": parsed["classification"]["layout"],
            "objects_count": len(parsed["objects"]),
        })

    # 1. 템플릿 생성 (MongoDB)
    template_doc = {
        "name": template_name,
        "description": f"PPTX 파일에서 자동 생성 ({total_slides}개 슬라이드)",
        "background_image": template_bg,
        "is_published": False,
        "created_by": user_key,
        "created_at": now,
        "updated_at": now,
    }
    result = await db.templates.insert_one(template_doc)
    template_id = str(result.inserted_id)

    # 2. 슬라이드 생성 (MongoDB - 일괄 삽입)
    slide_docs = []
    for idx, parsed in enumerate(parsed_slides):
        slide_doc = {
            "template_id": template_id,
            "order": idx + 1,
            "objects": parsed["objects"],
            "slide_meta": parsed["slide_meta"],
            "background_image": parsed["background_image"],
            "created_at": now,
            "updated_at": now,
        }
        slide_docs.append(slide_doc)

    if slide_docs:
        await db.slides.insert_many(slide_docs)

    logger.info(
        f"PPTX 임포트 완료: template_id={template_id}, "
        f"slides={total_slides}, user={user_key}"
    )

    return {
        "template_id": template_id,
        "slides_count": total_slides,
        "classification": classifications,
    }
