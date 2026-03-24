"""PPTX 생성 서비스 - python-pptx 기반 슬라이드 렌더링"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from services.mongo_service import get_db
from bson import ObjectId
from PIL import Image as PILImage
import os
import uuid
from config import settings

# 캔버스 / PPTX 상수 (기본 16:9)
SLIDE_W_PX = 960
SLIDE_H_PX = 540
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5148000

# 슬라이드 사이즈별 상수 (1cm = 360000 EMU)
# A4:   27.5 x 19.05 cm
# 16:9: 25.4 x 14.3 cm
# 4:3:  25.4 x 19.05 cm
SLIDE_SIZES = {
    "16:9": {"px_w": 960, "px_h": 540, "emu_w": 9144000,  "emu_h": 5148000},
    "4:3":  {"px_w": 960, "px_h": 720, "emu_w": 9144000,  "emu_h": 6858000},
    "A4":   {"px_w": 960, "px_h": 665, "emu_w": 9900000,  "emu_h": 6858000},
}


def _get_slide_size(slide_size: str = "16:9") -> dict:
    """슬라이드 사이즈 상수 반환"""
    return SLIDE_SIZES.get(slide_size, SLIDE_SIZES["16:9"])


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """HEX 색상을 RGBColor로 변환"""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) < 6:
        hex_color = "000000"
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _get_alignment(align: str):
    """텍스트 정렬 매핑"""
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
        align, PP_ALIGN.LEFT
    )


def _px_to_emu(x, y, w, h, slide_size="16:9"):
    """캔버스 좌표(px) → PPTX EMU 변환"""
    sz = _get_slide_size(slide_size)
    return (
        int(x / sz["px_w"] * sz["emu_w"]),
        int(y / sz["px_h"] * sz["emu_h"]),
        int(w / sz["px_w"] * sz["emu_w"]),
        int(h / sz["px_h"] * sz["emu_h"]),
    )


def _apply_cover_crop(pic, img_path: str, target_w_emu: int, target_h_emu: int):
    """image_fit='cover' 적용: 이미지를 영역에 맞춰 채우고 초과 부분을 크롭"""
    try:
        with PILImage.open(img_path) as img:
            img_w, img_h = img.size
        target_ratio = target_w_emu / target_h_emu
        img_ratio = img_w / img_h

        # 크롭 비율 계산 (0~100000 범위, 1000분율이 아니라 100000분율)
        if img_ratio > target_ratio:
            # 이미지가 더 넓음 → 좌우 크롭
            crop_w = 1 - (target_ratio / img_ratio)
            crop_each = int(crop_w / 2 * 100000)
            pic.crop_left = crop_each
            pic.crop_right = crop_each
        else:
            # 이미지가 더 높음 → 상하 크롭
            crop_h = 1 - (img_ratio / target_ratio)
            crop_each = int(crop_h / 2 * 100000)
            pic.crop_top = crop_each
            pic.crop_bottom = crop_each
    except Exception:
        pass


async def _generate_pptx_from_custom_template(project_id: str, project: dict) -> str:
    """커스텀 PPTX 템플릿을 기반으로 PPTX 파일 생성 (클론-앤-필)

    원본 PPTX를 열어서 생성된 텍스트로 교체합니다.
    """
    db = get_db()
    custom_template_id = project.get("custom_template_id", "")
    if not custom_template_id:
        raise ValueError("커스텀 템플릿 ID가 없습니다")

    custom_tmpl = await db.custom_templates.find_one({"_id": ObjectId(custom_template_id)})
    if not custom_tmpl:
        custom_tmpl = await db.custom_templates.find_one({"project_id": project_id})
    if not custom_tmpl:
        raise ValueError("커스텀 템플릿을 찾을 수 없습니다")

    file_path = custom_tmpl.get("file_path", "")
    abs_path = os.path.join(".", file_path.lstrip("/"))
    if not os.path.exists(abs_path):
        raise ValueError(f"커스텀 템플릿 파일을 찾을 수 없습니다: {file_path}")

    # 생성된 슬라이드 데이터 조회
    cursor = db.generated_slides.find({"project_id": project_id}).sort("order", 1)
    gen_slides = []
    async for s in cursor:
        gen_slides.append(s)

    if not gen_slides:
        raise ValueError("생성된 슬라이드가 없습니다")

    # 원본 PPTX 열기
    prs = Presentation(abs_path)
    original_slides = list(prs.slides)

    # 각 생성 슬라이드에 대해 원본 슬라이드의 텍스트 교체
    for gen_idx, gen_slide in enumerate(gen_slides):
        if gen_idx >= len(original_slides):
            break  # 원본보다 생성 슬라이드가 많으면 초과분은 무시

        pptx_slide = original_slides[gen_idx]
        gen_objects = gen_slide.get("objects", [])

        # generated_text가 있는 텍스트 오브젝트 수집
        text_replacements = []
        for obj in gen_objects:
            if obj.get("obj_type") == "text":
                generated_text = obj.get("generated_text", "")
                if generated_text:
                    text_replacements.append(generated_text)

        # 원본 슬라이드의 텍스트 프레임에 순서대로 교체
        text_idx = 0
        for shape in pptx_slide.shapes:
            if text_idx >= len(text_replacements):
                break
            if shape.has_text_frame:
                text_content = shape.text_frame.text.strip()
                if text_content:
                    # 첫 번째 paragraph의 run에 텍스트 교체
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        if p.runs:
                            p.runs[0].text = text_replacements[text_idx]
                        else:
                            run = p.add_run()
                            run.text = text_replacements[text_idx]
                    text_idx += 1

    # 파일 저장
    output_dir = os.path.join(settings.UPLOAD_DIR, "generated")
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)

    return f"/uploads/generated/{filename}"


async def generate_pptx(project_id: str) -> str:
    """생성된 슬라이드 데이터를 기반으로 PPTX 파일 생성"""
    db = get_db()

    project = await db.projects.find_one({"_id": ObjectId(project_id)})
    if not project:
        raise ValueError("프로젝트를 찾을 수 없습니다")

    # 커스텀 템플릿이 있으면 클론-앤-필 방식 사용
    if project.get("custom_template_id"):
        try:
            return await _generate_pptx_from_custom_template(project_id, project)
        except Exception as e:
            import logging
            logging.getLogger(__name__).warning(f"커스텀 템플릿 PPTX 생성 실패, 기본 방식으로 폴백: {e}")

    cursor = db.generated_slides.find({"project_id": project_id}).sort("order", 1)
    gen_slides = []
    async for s in cursor:
        gen_slides.append(s)

    if not gen_slides:
        raise ValueError("생성된 슬라이드가 없습니다")

    template = None
    if project.get("template_id"):
        template = await db.templates.find_one({"_id": ObjectId(project["template_id"])})

    prs = Presentation()
    _slide_size = template.get("slide_size", "16:9") if template else "16:9"
    slide_sz = _get_slide_size(_slide_size)
    prs.slide_width = Emu(slide_sz["emu_w"])
    prs.slide_height = Emu(slide_sz["emu_h"])

    for gen_slide in gen_slides:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 배경 이미지
        bg_image = gen_slide.get("background_image") or (
            template.get("background_image") if template else None
        )
        if bg_image:
            bg_path = os.path.join(".", bg_image.lstrip("/"))
            if os.path.exists(bg_path):
                try:
                    slide.shapes.add_picture(
                        bg_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height,
                    )
                except Exception:
                    pass

        # 오브젝트 렌더링 (z_index 순서로 정렬, items 소진된 초과 subtitle/description 제거)
        objects_sorted = sorted(gen_slide.get("objects", []), key=lambda o: o.get("z_index", 10))
        items = gen_slide.get("items", [])
        sub_idx = 0
        desc_idx = 0

        # 초과 subtitle/description의 Y 좌표 수집 → 같은 행의 number/shape도 제거
        removed_y = set()
        if items:
            _sub_i, _desc_i = 0, 0
            for o in objects_sorted:
                r = o.get("role", "") or o.get("_auto_role", "")
                ar = o.get("role_auto", False)
                if r == "subtitle" and not ar:
                    if _sub_i >= len(items):
                        removed_y.add(round(o.get("y", 0)))
                    _sub_i += 1
                elif r == "description" and not ar:
                    if _desc_i >= len(items):
                        removed_y.add(round(o.get("y", 0)))
                    _desc_i += 1

        for obj in objects_sorted:
            role = obj.get("role", "") or obj.get("_auto_role", "")
            is_auto_role = obj.get("role_auto", False)
            if items and role == "subtitle" and not is_auto_role:
                if sub_idx >= len(items):
                    continue
                sub_idx += 1
            elif items and role == "description" and not is_auto_role:
                if desc_idx >= len(items):
                    continue
                desc_idx += 1

            # 제거된 행의 number/shape도 함께 제거
            if removed_y and (obj.get("obj_type") == "shape" or role == "number"):
                obj_y = round(obj.get("y", 0))
                if any(abs(obj_y - ry) <= 15 for ry in removed_y):
                    continue

            xe, ye, we, he = _px_to_emu(obj["x"], obj["y"], obj["width"], obj["height"], _slide_size)

            if obj["obj_type"] == "image" and obj.get("image_url"):
                img_path = os.path.join(".", obj["image_url"].lstrip("/"))
                if os.path.exists(img_path):
                    try:
                        pic = slide.shapes.add_picture(
                            img_path, Emu(xe), Emu(ye), Emu(we), Emu(he),
                        )
                        if obj.get("image_fit") == "cover":
                            _apply_cover_crop(pic, img_path, we, he)
                    except Exception:
                        pass

            elif obj["obj_type"] == "text":
                _render_text(slide, obj, xe, ye, we, he)

            elif obj["obj_type"] == "shape":
                _render_shape(slide, obj, xe, ye, we, he)

            elif obj["obj_type"] == "table":
                _render_table(slide, obj, xe, ye, we, he)

            elif obj["obj_type"] == "chart":
                _render_chart(slide, obj, xe, ye, we, he)

    # 파일 저장
    output_dir = os.path.join(settings.UPLOAD_DIR, "generated")
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)

    return f"/uploads/generated/{filename}"


# ============ 텍스트 렌더링 ============

def _render_text(slide, obj: dict, x: int, y: int, w: int, h: int):
    """텍스트 오브젝트를 PPTX 슬라이드에 렌더링"""
    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    role = obj.get("role", "") or obj.get("_auto_role", "")
    if role == "description":
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    text_content = obj.get("generated_text") or obj.get("text_content", "")
    style = obj.get("text_style", {})

    p = tf.paragraphs[0]
    p.alignment = _get_alignment(style.get("align", "left"))

    run = p.add_run()
    run.text = text_content
    run.font.size = Pt(style.get("font_size", 16))
    run.font.bold = style.get("bold", False)
    run.font.italic = style.get("italic", False)
    run.font.color.rgb = _hex_to_rgb(style.get("color", "#000000"))

    if style.get("font_family"):
        run.font.name = style["font_family"]


# ============ 도형 렌더링 ============

def _render_shape(slide, obj: dict, x: int, y: int, w: int, h: int):
    """도형 오브젝트를 PPTX 슬라이드에 렌더링"""
    style = obj.get("shape_style", {})
    shape_type = style.get("shape_type", "rectangle")

    if shape_type in ("line", "arrow"):
        _render_line_shape(slide, style, x, y, w, h)
    else:
        _render_block_shape(slide, style, shape_type, x, y, w, h)


def _render_block_shape(slide, style: dict, shape_type: str, x: int, y: int, w: int, h: int):
    """블록 도형 렌더링 - 직사각형, 기본 도형, 블록 화살표, 수학, 별/현수막, 설명선"""
    shape_map = {
        # 직사각형
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "snip_1_rect": MSO_SHAPE.SNIP_1_RECTANGLE,
        "snip_2_diag_rect": MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
        "round_1_rect": MSO_SHAPE.ROUND_1_RECTANGLE,
        "round_2_diag_rect": MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
        # 기본 도형
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
        "diamond": MSO_SHAPE.DIAMOND,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "heptagon": MSO_SHAPE.HEPTAGON,
        "octagon": MSO_SHAPE.OCTAGON,
        "decagon": MSO_SHAPE.DECAGON,
        "dodecagon": MSO_SHAPE.DODECAGON,
        "cross": MSO_SHAPE.CROSS,
        "donut": MSO_SHAPE.DONUT,
        "no_smoking": MSO_SHAPE.NO_SYMBOL,
        "block_arc": MSO_SHAPE.BLOCK_ARC,
        "heart": MSO_SHAPE.HEART,
        "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
        "sun": MSO_SHAPE.SUN,
        "moon": MSO_SHAPE.MOON,
        "cloud": MSO_SHAPE.CLOUD,
        "smiley_face": MSO_SHAPE.SMILEY_FACE,
        "folded_corner": MSO_SHAPE.FOLDED_CORNER,
        "frame": MSO_SHAPE.FRAME,
        "teardrop": MSO_SHAPE.TEAR,
        "plaque": MSO_SHAPE.PLAQUE,
        "brace_pair": MSO_SHAPE.DOUBLE_BRACE,
        "bracket_pair": MSO_SHAPE.DOUBLE_BRACKET,
        # 블록 화살표
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "up_arrow": MSO_SHAPE.UP_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
        "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
        "up_down_arrow": MSO_SHAPE.UP_DOWN_ARROW,
        "quad_arrow": MSO_SHAPE.QUAD_ARROW,
        "notched_right_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
        "chevron": MSO_SHAPE.CHEVRON,
        "home_plate": MSO_SHAPE.PENTAGON,
        "striped_right_arrow": MSO_SHAPE.STRIPED_RIGHT_ARROW,
        "bent_arrow": MSO_SHAPE.BENT_ARROW,
        "u_turn_arrow": MSO_SHAPE.U_TURN_ARROW,
        "circular_arrow": MSO_SHAPE.CIRCULAR_ARROW,
        "curved_right_arrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
        "left_up_arrow": MSO_SHAPE.LEFT_UP_ARROW,
        "bent_up_arrow": MSO_SHAPE.BENT_UP_ARROW,
        # 수학
        "math_plus": MSO_SHAPE.MATH_PLUS,
        "math_minus": MSO_SHAPE.MATH_MINUS,
        "math_multiply": MSO_SHAPE.MATH_MULTIPLY,
        "math_divide": MSO_SHAPE.MATH_DIVIDE,
        "math_equal": MSO_SHAPE.MATH_EQUAL,
        "math_not_equal": MSO_SHAPE.MATH_NOT_EQUAL,
        # 별 및 현수막
        "star_4_point": MSO_SHAPE.STAR_4_POINT,
        "star_5_point": MSO_SHAPE.STAR_5_POINT,
        "star_6_point": MSO_SHAPE.STAR_6_POINT,
        "star_8_point": MSO_SHAPE.STAR_8_POINT,
        "star_10_point": MSO_SHAPE.STAR_10_POINT,
        "star_12_point": MSO_SHAPE.STAR_12_POINT,
        "star_16_point": MSO_SHAPE.STAR_16_POINT,
        "star_24_point": MSO_SHAPE.STAR_24_POINT,
        "star_32_point": MSO_SHAPE.STAR_32_POINT,
        "explosion_1": MSO_SHAPE.EXPLOSION1,
        "explosion_2": MSO_SHAPE.EXPLOSION2,
        "wave": MSO_SHAPE.WAVE,
        "double_wave": MSO_SHAPE.DOUBLE_WAVE,
        "ribbon": MSO_SHAPE.DOWN_RIBBON,
        # 설명선
        "wedge_rect_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
        "wedge_round_rect_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "wedge_ellipse_callout": MSO_SHAPE.OVAL_CALLOUT,
        "cloud_callout": MSO_SHAPE.CLOUD_CALLOUT,
        "border_callout_1": MSO_SHAPE.LINE_CALLOUT_1,
        "border_callout_2": MSO_SHAPE.LINE_CALLOUT_2,
        "border_callout_3": MSO_SHAPE.LINE_CALLOUT_3,
    }
    shape_enum = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_enum, Emu(x), Emu(y), Emu(w), Emu(h))

    # 채우기
    fill_opacity = style.get("fill_opacity", 1.0)
    fill_color = style.get("fill_color", "#4A90D9")

    if fill_opacity == 0 or fill_color in ("transparent", "none"):
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        if fill_opacity < 1.0:
            _set_fill_opacity(shape, fill_opacity)

    # 둥근 사각형 커스텀 border_radius
    if shape_type == "rounded_rectangle" and style.get("border_radius"):
        _set_corner_radius(shape, style["border_radius"], w, h)

    # 테두리
    stroke_width = style.get("stroke_width", 2)
    if stroke_width > 0:
        shape.line.color.rgb = _hex_to_rgb(style.get("stroke_color", "#333333"))
        shape.line.width = Pt(stroke_width)
        _set_line_dash(shape, style.get("stroke_dash", "solid"))
    else:
        shape.line.fill.background()


def _render_line_shape(slide, style: dict, x: int, y: int, w: int, h: int):
    """라인/화살표 도형 렌더링"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(max(h, Pt(1).emu))
    )

    shape.fill.background()

    stroke_width = style.get("stroke_width", 2)
    shape.line.color.rgb = _hex_to_rgb(style.get("stroke_color", "#333333"))
    shape.line.width = Pt(stroke_width)
    _set_line_dash(shape, style.get("stroke_dash", "solid"))

    if style.get("shape_type") == "arrow":
        try:
            ln_elem = shape.line._ln
            tail_end = ln_elem.makeelement(qn("a:tailEnd"), {})
            tail_end.set("type", "triangle")
            tail_end.set("w", "med")
            tail_end.set("len", "med")
            ln_elem.append(tail_end)

            if style.get("arrow_head") == "both":
                head_end = ln_elem.makeelement(qn("a:headEnd"), {})
                head_end.set("type", "triangle")
                head_end.set("w", "med")
                head_end.set("len", "med")
                ln_elem.append(head_end)
        except Exception:
            pass


# ============ 테이블 렌더링 ============

def _render_table(slide, obj: dict, x: int, y: int, w: int, h: int):
    """테이블 오브젝트를 PPTX 슬라이드에 렌더링"""
    style = obj.get("table_style", {})
    rows = style.get("rows", 3)
    cols = style.get("cols", 3)
    data = style.get("data", [])
    merged_cells = style.get("merged_cells", [])
    cell_styles = style.get("cell_styles", {})

    if rows < 1 or cols < 1:
        return

    table_shape = slide.shapes.add_table(rows, cols, Emu(x), Emu(y), Emu(w), Emu(h))
    table = table_shape.table

    font_family = style.get("font_family", "Arial")
    font_size = style.get("font_size", 11)
    header_row = style.get("header_row", True)
    header_bg = style.get("header_bg_color", "#4472C4")
    header_text = style.get("header_text_color", "#FFFFFF")
    banded_rows = style.get("banded_rows", False)
    border_color = style.get("border_color", "#BFBFBF")
    border_width = style.get("border_width", 1)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell_text = ""
            if r < len(data) and c < len(data[r]):
                cell_text = str(data[r][c])

            # Clear default paragraph and set text
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(font_size)
            if font_family:
                run.font.name = font_family

            # Per-cell style overrides
            cell_key = f"{r}_{c}"
            cs = cell_styles.get(cell_key, {})

            # Header row styling
            if header_row and r == 0:
                bg = cs.get("bg_color", header_bg)
                tc = cs.get("text_color", header_text)
                _set_cell_fill(cell, bg)
                run.font.color.rgb = _hex_to_rgb(tc)
                run.font.bold = cs.get("bold", True)
            else:
                if cs.get("bg_color"):
                    _set_cell_fill(cell, cs["bg_color"])
                elif banded_rows and r % 2 == 0 and (not header_row or r > 0):
                    _set_cell_fill(cell, "#D9E2F3")

                if cs.get("text_color"):
                    run.font.color.rgb = _hex_to_rgb(cs["text_color"])
                if cs.get("bold"):
                    run.font.bold = True

            # Text alignment
            align = cs.get("text_align", "left")
            p.alignment = _get_alignment(align)

            # Cell borders
            _set_cell_borders(cell, border_color, border_width)

    # Merge cells
    for merge in merged_cells:
        try:
            sr, sc = merge.get("start_row", 0), merge.get("start_col", 0)
            er, ec = merge.get("end_row", 0), merge.get("end_col", 0)
            if sr < rows and sc < cols and er < rows and ec < cols:
                table.cell(sr, sc).merge(table.cell(er, ec))
        except Exception:
            pass


def _set_cell_fill(cell, hex_color: str):
    """테이블 셀 배경색 설정"""
    try:
        tc_pr = cell._tc.get_or_add_tcPr()
        solid_fill = tc_pr.makeelement(qn("a:solidFill"), {})
        srgb_clr = solid_fill.makeelement(qn("a:srgbClr"), {"val": hex_color.lstrip("#")})
        solid_fill.append(srgb_clr)
        # Remove existing fill
        for existing in tc_pr.findall(qn("a:solidFill")):
            tc_pr.remove(existing)
        tc_pr.append(solid_fill)
    except Exception:
        pass


def _set_cell_borders(cell, border_color: str, border_width: int):
    """테이블 셀 테두리 설정"""
    try:
        tc_pr = cell._tc.get_or_add_tcPr()
        border_width_emu = Pt(border_width).emu
        color_val = border_color.lstrip("#")

        for border_name in ["lnL", "lnR", "lnT", "lnB"]:
            # Remove existing
            for existing in tc_pr.findall(qn(f"a:{border_name}")):
                tc_pr.remove(existing)
            ln = tc_pr.makeelement(qn(f"a:{border_name}"), {"w": str(border_width_emu)})
            solid_fill = ln.makeelement(qn("a:solidFill"), {})
            srgb = solid_fill.makeelement(qn("a:srgbClr"), {"val": color_val})
            solid_fill.append(srgb)
            ln.append(solid_fill)
            tc_pr.append(ln)
    except Exception:
        pass


# ============ 차트 렌더링 ============

def _render_chart(slide, obj: dict, x: int, y: int, w: int, h: int):
    """차트 오브젝트를 PPTX 슬라이드에 렌더링"""
    style = obj.get("chart_style", {})
    chart_type = style.get("chart_type", "bar")
    chart_data_raw = style.get("chart_data", {})
    labels = chart_data_raw.get("labels", [])
    datasets = chart_data_raw.get("datasets", [])

    if not labels or not datasets:
        return

    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "radar": XL_CHART_TYPE.RADAR,
    }

    chart_data = CategoryChartData()
    chart_data.categories = labels
    for ds in datasets:
        series_data = ds.get("data", [])
        # Ensure data length matches categories
        while len(series_data) < len(labels):
            series_data.append(0)
        chart_data.add_series(ds.get("label", "Series"), tuple(series_data[:len(labels)]))

    xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    try:
        chart_frame = slide.shapes.add_chart(xl_type, Emu(x), Emu(y), Emu(w), Emu(h), chart_data)
        chart = chart_frame.chart

        # Legend
        chart.has_legend = style.get("show_legend", True)
        if chart.has_legend:
            chart.legend.include_in_layout = False

        # Title
        chart_title = style.get("title", "")
        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_title
        else:
            chart.has_title = False

        # Apply color scheme to series
        color_scheme = style.get("color_scheme", ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000", "#5B9BD5", "#70AD47"])
        try:
            plot = chart.plots[0]
            for i, series in enumerate(plot.series):
                color = color_scheme[i % len(color_scheme)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = _hex_to_rgb(color)
        except Exception:
            pass

    except Exception:
        pass


# ============ 유틸리티 ============

def _set_fill_opacity(shape, opacity: float):
    """도형 채우기 투명도 설정 (XML 직접 조작)"""
    try:
        solid_fill = shape.fill._fill.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = srgb.makeelement(qn("a:alpha"), {})
                alpha.set("val", str(int(opacity * 100000)))
                srgb.append(alpha)
    except Exception:
        pass


def _set_corner_radius(shape, border_radius_px: int, w_emu: int, h_emu: int):
    """둥근 사각형 모서리 반경 커스텀 설정"""
    try:
        min_dim_emu = min(w_emu, h_emu)
        if min_dim_emu <= 0:
            return
        radius_emu = int(border_radius_px / SLIDE_W_PX * SLIDE_W_EMU)
        adj_val = int(radius_emu / (min_dim_emu / 2) * 50000)
        adj_val = max(0, min(adj_val, 50000))

        sp = shape._element
        prst_geom = sp.find(qn("a:prstGeom"))
        if prst_geom is not None:
            av_lst = prst_geom.find(qn("a:avLst"))
            if av_lst is None:
                av_lst = prst_geom.makeelement(qn("a:avLst"), {})
                prst_geom.append(av_lst)
            for existing in av_lst.findall(qn("a:gd")):
                av_lst.remove(existing)
            gd = av_lst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {adj_val}"})
            av_lst.append(gd)
    except Exception:
        pass


def _set_line_dash(shape, dash_style: str):
    """테두리 선 스타일 설정"""
    try:
        if dash_style == "dashed":
            shape.line.dash_style = 4  # DASH
        elif dash_style == "dotted":
            shape.line.dash_style = 3  # ROUND_DOT
    except Exception:
        pass
